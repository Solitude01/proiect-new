import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- 全局配置 ---
# 修改文件名以体现新尺寸
FILENAME = "2025年党员攻坚项目汇报_AI侦测项目4.pptx"
AUTHOR = "吴佳宇"

# 配色方案 (SCC Corporate Identity) - 保持不变
C_BLUE = RGBColor(0, 51, 102)        # #003366 主色
C_GOLD = RGBColor(184, 134, 11)      # #B8860B 辅助色
C_GRAY_TEXT = RGBColor(51, 51, 51)   # #333333 正文
C_LIGHT_BG = RGBColor(248, 249, 250) # #F8F9FA 卡片背景
C_RED = RGBColor(211, 47, 47)        # #D32F2F 警告
C_GREEN = RGBColor(40, 167, 69)      # #28A745 完成
C_WHITE = RGBColor(255, 255, 255)    # #FFFFFF
C_BORDER = RGBColor(224, 224, 224)   # #E0E0E0

def create_presentation():
    # 1. 初始化 PPT画布
    # 【关键修改】设置为标准 16:9 尺寸 (10 x 5.625 英寸)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 设置核心属性
    prs.core_properties.author = AUTHOR
    prs.core_properties.title = "AI侦测辅助监控项目汇报"

    # --- 核心辅助函数 (保持不变) ---
    def add_textbox(slide, text, x, y, w, h, font_size=12, color=C_GRAY_TEXT, bold=False, align=None, bg_color=None, border_color=None):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        if bg_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if align:
            p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = 'Microsoft YaHei'
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        return shape

    def add_card(slide, title, content, x, y, w, h, highlight=False):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_LIGHT_BG
        shape.line.color.rgb = C_GOLD if highlight else C_BLUE
        shape.line.width = Pt(1.5 if highlight else 0.75)
        
        title_box = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+0.1), Inches(w-0.2), Inches(0.4))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = title
        tp.font.bold = True
        tp.font.size = Pt(13)
        tp.font.name = 'Microsoft YaHei'
        tp.font.color.rgb = C_BLUE
        
        content_box = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+0.5), Inches(w-0.2), Inches(h-0.6))
        content_box.text_frame.word_wrap = True
        cp = content_box.text_frame.paragraphs[0]
        cp.text = content
        cp.font.size = Pt(11)
        cp.font.name = 'Microsoft YaHei'
        cp.font.color.rgb = C_GRAY_TEXT

    # --- 核心辅助函数：应用母版样式 (所有坐标已缩放 x0.75) ---
    def apply_master(slide, title_text, sub_text=None, page_num=None):
        # 1. 左上角 Logo (原 0.4, 0.4)
        logo = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(3), Inches(0.5))
        logo.text_frame.text = "SCC深南电路"
        logo.text_frame.paragraphs[0].font.bold = True
        logo.text_frame.paragraphs[0].font.color.rgb = C_GOLD
        logo.text_frame.paragraphs[0].font.size = Pt(14)

        # 2. 标题与分割线 (原 y=0.8, 0.95)
        add_textbox(slide, title_text, 0.3, 0.6, 7, 0.6, 22, C_BLUE, True)
        if sub_text:
             add_textbox(slide, f"|  {sub_text}", 4.0, 0.7, 5.5, 0.4, 14, C_GRAY_TEXT)
        
        # 分割线 (原 y=1.5, w=12.53)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.125), Inches(9.4), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = C_GOLD
        line.line.fill.background()

        # 3. 页脚 (原 y=7.1)
        footer_y = 5.325
        footer = slide.shapes.add_textbox(Inches(0.3), Inches(footer_y), Inches(5), Inches(0.3))
        footer.text_frame.text = "2025年党员攻坚项目 | 南通深南党支部"
        footer.text_frame.paragraphs[0].font.size = Pt(9)
        footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(170, 170, 170)

        if page_num:
            # (原 x=12.5)
            pg = slide.shapes.add_textbox(Inches(9.2), Inches(footer_y), Inches(0.5), Inches(0.3))
            pg.text_frame.text = str(page_num)
            pg.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            pg.text_frame.paragraphs[0].font.size = Pt(9)
            pg.text_frame.paragraphs[0].font.color.rgb = RGBColor(170, 170, 170)

    # --- 核心辅助函数：添加演讲稿框 (所有坐标已缩放 x0.75) ---
    def add_script(slide, script_text):
        # (原 y=5.8, w=12.53, h=1.1)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(4.35), Inches(9.4), Inches(0.85))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(245, 245, 245)
        bg.line.color.rgb = RGBColor(204, 204, 204)
        bg.line.dash_style = 1 
        
        # (原 y=5.9)
        add_textbox(slide, "演讲关键句：", 0.4, 4.425, 2, 0.3, 11, C_GOLD, True)
        # (原 y=6.2, w=12.3)
        content = slide.shapes.add_textbox(Inches(0.4), Inches(4.65), Inches(9.2), Inches(0.5))
        content.text_frame.word_wrap = True
        p = content.text_frame.paragraphs[0]
        p.text = script_text
        p.font.italic = True
        p.font.size = Pt(10.5) # 稍微调小字体
        p.font.color.rgb = RGBColor(85, 85, 85)

    # ==================== 封面页 (Cover) ====================
    slide_cover = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 封面背景 (原 y=5.5, h=2)
    shape_bg = slide_cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.125), Inches(10), Inches(1.5))
    shape_bg.fill.solid(); shape_bg.fill.fore_color.rgb = C_BLUE; shape_bg.line.fill.background()
    
    # Logo (原 0.5, 0.5)
    logo_cover = slide_cover.shapes.add_textbox(Inches(0.375), Inches(0.375), Inches(3), Inches(0.5))
    logo_cover.text_frame.text = "SCC深南电路"
    logo_cover.text_frame.paragraphs[0].font.bold = True
    logo_cover.text_frame.paragraphs[0].font.color.rgb = C_GOLD
    logo_cover.text_frame.paragraphs[0].font.size = Pt(18)

    # 标题 (原 y=2.5, 3.5) - 字号微调
    add_textbox(slide_cover, "AI侦测辅助监控项目", 0.75, 1.875, 8, 1.2, 40, C_BLUE, True)
    add_textbox(slide_cover, "2025年度党员攻坚项目汇报", 0.75, 2.7, 6, 0.6, 20, C_GRAY_TEXT)
    
    # 金色分割线 (原 y=4.3)
    line_cover = slide_cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(3.3), Inches(1.125), Inches(0.08))
    line_cover.fill.solid(); line_cover.fill.fore_color.rgb = C_GOLD; line_cover.line.fill.background()

    # 底部信息 (原 y=6.0, 6.5)
    add_textbox(slide_cover, "汇报支部：南通深南党支部", 0.75, 4.5, 5, 0.4, 14, C_WHITE, True)
    add_textbox(slide_cover, "项目负责人：吴佳宇", 0.75, 4.9, 5, 0.4, 14, C_WHITE, False)

    # 装饰圆环 (原 x=9, y=1.5, size=5)
    circle = slide_cover.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.75), Inches(1.125), Inches(3.75), Inches(3.75))
    circle.fill.background(); circle.line.color.rgb = RGBColor(240, 240, 240); circle.line.width = Pt(15)

    # ==================== P1: 政治引领与顶层设计 ====================
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s1, "01 政治引领与顶层设计", "筑牢战斗堡垒，聚焦数字化瓶颈", 1)

    # 上半部卡片 (原 y=1.8, w=6.0, h=1.8)
    card_y = 1.35
    card_w = 4.5
    card_h = 1.35
    
    # 左卡片 (原 x=0.5)
    add_card(s1, "🚩 政治站位与方针", 
             "● 落实上级党委战略部署，确立“党建+提质增效”总方针。\n● 将本项目定为支部级攻坚任务，服务高质量发展。", 
             0.375, card_y, card_w, card_h)
    
    # 右卡片 (原 x=6.8)
    add_card(s1, "⚠️ 核心挑战 (痛点)", 
             "● 传统人工点检模式效率低下，有效性仅20%。\n● 严重制约公司数字化转型深化的进程。\n● 属于数字化推进中的“卡脖子”难题。", 
             5.1, card_y, card_w, card_h, highlight=True)

    # 下半部目标区域 (原 y=3.9, w=12.3, h=1.5)
    goal_bg_y = 2.925
    goal_bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.375), Inches(goal_bg_y), Inches(9.225), Inches(1.125))
    goal_bg.fill.solid(); goal_bg.fill.fore_color.rgb = C_LIGHT_BG
    goal_bg.line.color.rgb = C_BLUE
    
    # 标题 (原 y=4.0)
    add_textbox(s1, "🎯 顶层设计与关键目标", 0.45, 3.0, 4, 0.3, 12, C_BLUE, True)
    
    # 目标数据 (原 y=4.4, 5.1)
    data_y = 3.3
    label_y = 3.75
    # 目标 1 (原 x=1.0)
    add_textbox(s1, "90%+", 0.75, data_y, 1.5, 0.6, 28, C_BLUE, True, PP_ALIGN.CENTER)
    add_textbox(s1, "点检准确率", 0.75, label_y, 1.5, 0.3, 10, C_GRAY_TEXT, False, PP_ALIGN.CENTER)
    
    # 目标 2 (原 x=4.0)
    add_textbox(s1, "AI对接", 3.0, data_y, 1.5, 0.6, 28, C_BLUE, True, PP_ALIGN.CENTER)
    add_textbox(s1, "系统无缝集成", 3.0, label_y, 1.5, 0.3, 10, C_GRAY_TEXT, False, PP_ALIGN.CENTER)
    
    # 描述文本 (原 x=7.0, y=4.4)
    add_textbox(s1, "通过AI技术革新，实现从“人工抽检”向“智能全检”的跨越，彻底解决异常发现滞后问题。", 
                5.25, 3.3, 4.0, 0.6, 11, C_GRAY_TEXT)

    add_script(s1, "“面对人工点检有效性仅20%的瓶颈，支部主动领题，将AI侦测项目确立为年度攻坚任务，以‘党建+提质增效’为指引，誓要攻克这一数字化转型的卡脖子难题。”")

    # ==================== P2: 组织攻坚与党群联动 ====================
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s2, "02 组织攻坚与党群联动", "党员领题突击，群众聚力攻坚", 2)

    # 架构标题 (原 y=1.6)
    add_textbox(s2, "“1+N”党群联动攻坚体系", 0.375, 1.2, 5, 0.4, 15, C_BLUE, True)

    # 1. 战斗堡垒 (顶层) (原 x=4.5, y=2.2, w=4.3, h=0.8)
    fortress = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.375), Inches(1.65), Inches(3.225), Inches(0.6))
    fortress.fill.solid(); fortress.fill.fore_color.rgb = C_BLUE
    add_textbox(s2, "🛡️ 南通深南党支部：“AI数字赋能”突击队", 3.45, 1.75, 3.075, 0.4, 11, C_WHITE, True, PP_ALIGN.CENTER)

    # 连接线 (原 x=6.6, y=3.0, h=0.4)
    conn = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.95), Inches(2.25), Inches(0.04), Inches(0.3))
    conn.fill.solid(); conn.fill.fore_color.rgb = C_GOLD

    # 下方三个框体通用参数 (原 y=3.4, w=3.8, h=2.2)
    box_y = 2.55
    box_w = 2.85
    box_h = 1.65

    # 左：负责人 (原 x=0.5)
    box_l_x = 0.375
    box_l = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(box_l_x), Inches(box_y), Inches(box_w), Inches(box_h))
    box_l.fill.solid(); box_l.fill.fore_color.rgb = C_LIGHT_BG; box_l.line.color.rgb = C_BLUE
    add_textbox(s2, "👤 项目负责人 (党员)", box_l_x+0.1, box_y+0.1, box_w-0.2, 0.3, 11, C_BLUE, True)
    add_textbox(s2, "吴佳宇", box_l_x+0.1, box_y+0.4, box_w-0.2, 0.3, 13, C_GRAY_TEXT, True)
    add_textbox(s2, "职责：资源调配 / 风险控制\n战绩：在素材采集困难时确保方向不偏航。", box_l_x+0.1, box_y+0.8, box_w-0.2, 0.7, 10, C_GRAY_TEXT)

    # 中：技术骨干 (原 x=4.7)
    box_m_x = 3.525
    box_m = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(box_m_x), Inches(box_y), Inches(box_w), Inches(box_h))
    box_m.fill.solid(); box_m.fill.fore_color.rgb = C_LIGHT_BG; box_m.line.color.rgb = C_BLUE
    add_textbox(s2, "🔧 技术攻坚 (党员骨干)", box_m_x+0.1, box_y+0.1, box_w-0.2, 0.3, 11, C_BLUE, True)
    add_textbox(s2, "潘秦", box_m_x+0.1, box_y+0.4, box_w-0.2, 0.3, 13, C_GRAY_TEXT, True)
    add_textbox(s2, "职责：算法预研 / 模型落地\n战绩：攻克技术壁垒，模型有效性达99%。", box_m_x+0.1, box_y+0.8, box_w-0.2, 0.7, 10, C_GRAY_TEXT)

    # 右：群众联动 (原 x=8.9)
    box_r_x = 6.675
    box_r = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(box_r_x), Inches(box_y), Inches(box_w), Inches(box_h))
    box_r.fill.solid(); box_r.fill.fore_color.rgb = C_LIGHT_BG; box_r.line.color.rgb = C_GOLD
    add_textbox(s2, "🤝 群众联动 (1+N)", box_r_x+0.1, box_y+0.1, box_w-0.2, 0.3, 11, C_GOLD, True)
    add_textbox(s2, "张旭阳(预备) / 李鹏飞 / 丁灵", box_r_x+0.1, box_y+0.4, box_w-0.2, 0.3, 11, C_GRAY_TEXT, True)
    add_textbox(s2, "职责：配合点检 / 素材采集 / 系统开发\n成效：形成联合攻关合力。", box_r_x+0.1, box_y+0.8, box_w-0.2, 0.7, 10, C_GRAY_TEXT)

    add_script(s2, "“吴佳宇同志把控全局，潘秦同志攻克99%有效性的技术难关，同时带动张旭阳、李鹏飞等群众骨干参与，充分体现了‘党员带头、全员攻坚’的战斗堡垒作用。”")

    # ==================== P3: 重点项目全景推进表 ====================
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s3, "03 重点项目全景推进表", "挂图作战，节点管控，实绩说话", 3)

    # 表格参数缩放
    # 表头 y=1.8 -> 1.35, h=0.5 -> 0.375
    header_y = 1.35
    header_h = 0.375
    # 列宽缩放 (1.5->1.125, 4.5->3.375, 2.0->1.5, 4.0->3.0)
    cols = [1.125, 3.375, 1.5, 3.0]
    headers = ["节点时间", "关键里程碑", "完成状态", "备注/难点"]
    
    cx = 0.375
    for i, h in enumerate(headers):
        shape = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(header_y), Inches(cols[i]), Inches(header_h))
        shape.fill.solid(); shape.fill.fore_color.rgb = C_BLUE
        add_textbox(s3, h, cx, header_y+0.05, cols[i], 0.3, 10, C_WHITE, True, PP_ALIGN.CENTER)
        cx += cols[i] + 0.075

    # 数据行参数缩放
    data = [
        ("3月15日", "现状数据确认及分析", "✅ 已完成", ""),
        ("4月10日", "硬件调整及优化", "✅ 已完成", "高效完成"),
        ("6月15日", "软件开发级上线应用", "✅ 已完成", "核心节点"),
        ("9月20日", "互通模型/系统现场测试", "✅ 已完成", ""),
        ("推进中", "模型实现 (有效性99%)", "⏳ 进行中", "难点：应对工艺流程调整风险"),
        ("11月25日", "成果移交与规范化培训", "📅 计划中", "下一步里程碑")
    ]

    # 数据行起始 y=2.4 -> 1.8, 行高 0.5 -> 0.375, 行间距 0.6 -> 0.45
    cy = 1.8
    row_h = 0.375
    row_gap = 0.45

    for date, task, status, note in data:
        cx = 0.375
        row_colors = [C_GRAY_TEXT, C_GRAY_TEXT, C_GREEN if "已完成" in status else (C_BLUE if "进行中" in status else C_GRAY_TEXT), C_RED if "难点" in note else C_GRAY_TEXT]
        vals = [date, task, status, note]
        
        for i, val in enumerate(vals):
            bg_rect = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(cy), Inches(cols[i]), Inches(row_h))
            bg_rect.fill.solid(); bg_rect.fill.fore_color.rgb = C_LIGHT_BG
            bg_rect.line.color.rgb = C_BORDER
            # 字体调整为 10pt
            add_textbox(s3, val, cx, cy+0.03, cols[i], 0.3, 10, row_colors[i], False, PP_ALIGN.CENTER if i != 1 else PP_ALIGN.LEFT)
            cx += cols[i] + 0.075
        cy += row_gap

    add_script(s3, "“我们严格落实挂图作战。目前前四个关键节点均已高质量完成，当前正集中精力攻克模型应对工艺变化的难题，确保11月25日顺利完成项目移交。”")

    # ==================== P4: 赋能增效与数据实绩 ====================
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s4, "04 赋能增效与数据实绩", "将党建活力转化为发展动力", 4)

    # 左侧对比框 (原 x=0.5, y=1.8, w=5.5, h=3.5)
    box_ba_x = 0.375
    box_ba_y = 1.35
    box_ba_w = 4.125
    box_ba_h = 2.625
    box_ba = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(box_ba_x), Inches(box_ba_y), Inches(box_ba_w), Inches(box_ba_h))
    box_ba.fill.solid(); box_ba.fill.fore_color.rgb = C_LIGHT_BG; box_ba.line.color.rgb = C_BORDER

    add_textbox(s4, "📊 核心指标飞跃", box_ba_x+0.15, box_ba_y+0.15, 3, 0.3, 12, C_BLUE, True)

    # Before 数据 (原 x=0.8, y=2.6 / 3.5)
    data_x_l = 0.6
    data_y_val = 1.95
    data_y_lbl = 2.625
    add_textbox(s4, "20%", data_x_l, data_y_val, 1.5, 0.6, 28, C_GRAY_TEXT, True, PP_ALIGN.CENTER)
    add_textbox(s4, "Before (人工点检)", data_x_l, data_y_lbl, 1.5, 0.3, 10, C_GRAY_TEXT, False, PP_ALIGN.CENTER)

    # 箭头 (原 x=3.0, y=3.0, w=0.8, h=0.5)
    arr = s4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.25), Inches(2.25), Inches(0.6), Inches(0.375))
    arr.fill.solid(); arr.fill.fore_color.rgb = C_GOLD

    # After 数据 (原 x=3.8)
    data_x_r = 2.85
    add_textbox(s4, ">90%", data_x_r, data_y_val, 1.5, 0.6, 28, C_BLUE, True, PP_ALIGN.CENTER)
    add_textbox(s4, "After (AI侦测)", data_x_r, data_y_lbl, 1.5, 0.3, 10, C_BLUE, True, PP_ALIGN.CENTER)

    # 结论 (原 y=4.5)
    add_textbox(s4, "✅ 彻底解决异常发现滞后顽疾", box_ba_x, 3.4, box_ba_w, 0.4, 11, C_GREEN, True, PP_ALIGN.CENTER)

    # 右侧价值清单 (原 x=6.5, y=1.8)
    right_x = 4.875
    right_y_start = 1.35
    add_textbox(s4, "🚀 多维价值产出", right_x, right_y_start, 3, 0.4, 13, C_BLUE, True)
    
    vals = [
        ("技术革新", "成功预研并落地2个AI侦测辅助监控项目，实现技术突破。"),
        ("质量赋能", "有效弱化‘管理X’对产品质量的影响，增强监控时效。"),
        ("长效固化", "形成规范化操作流程，杜绝问题反弹，确保长期改善。")
    ]
    
    # (原 y start=2.4, increment=1.0)
    y_val = 1.8
    y_inc = 0.75
    for vt, vd in vals:
        add_textbox(s4, f"🔹 {vt}", right_x, y_val, 4.5, 0.3, 11, C_BLUE, True)
        add_textbox(s4, vd, right_x, y_val+0.3, 4.5, 0.45, 10, C_GRAY_TEXT)
        y_val += y_inc

    add_script(s4, "“项目成效显著：点检准确率实现质的飞跃，成功落地2个AI项目，切实将党建攻坚的活力转化为了提升产品质量的实际动力。”")

    # ==================== P5: 问题剖析与长效机制 ====================
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s5, "05 问题剖析与长效机制", "刀刃向内找差距，着眼长远建机制", 5)

    # 1. 反思框 (原 y=1.8, w=12.3, h=1.4)
    reflect_y = 1.35
    reflect_w = 9.225
    reflect_h = 1.05
    reflect_box = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.375), Inches(reflect_y), Inches(reflect_w), Inches(reflect_h))
    reflect_box.fill.solid(); reflect_box.fill.fore_color.rgb = RGBColor(255, 235, 238)
    reflect_box.line.color.rgb = C_RED
    reflect_box.line.dash_style = 1
    
    add_textbox(s5, "🛑 反思与差距", 0.5, reflect_y+0.1, 3, 0.3, 12, C_RED, True)
    add_textbox(s5, "● 风险警惕：‘管理X’随工艺变化调整可能导致模型失效。\n● 思想防范：持续防范技术与业务“两张皮”的思想残余。", 
                0.5, reflect_y+0.4, 8.5, 0.6, 10, C_GRAY_TEXT)

    # 2. PDCA机制 (原 y=3.5)
    pdca_y_title = 2.625
    add_textbox(s5, "🔄 后续重点与长效机制 (PDCA)", 0.375, pdca_y_title, 5, 0.3, 13, C_BLUE, True)
    
    # 步骤图参数缩放
    # (原 x start=0.8, y=4.1, w=2.5, h=1.4, increment=3.1)
    step_x = 0.6
    step_y = 3.1
    step_w = 1.875
    step_h = 1.05
    step_inc = 2.325
    
    steps = [
        ("成果固化", "11月25日前\n完成移交"),
        ("赋能培训", "实现系统在现场\n深度应用"),
        ("持续优化", "党建带团建\n模型持续迭代"),
        ("复制推广", "成功经验\n全公司推广")
    ]
    
    for i, (st, sd) in enumerate(steps):
        # 框
        s_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(step_x), Inches(step_y), Inches(step_w), Inches(step_h))
        s_box.fill.solid(); s_box.fill.fore_color.rgb = C_WHITE
        s_box.line.color.rgb = C_BLUE
        s_box.line.width = Pt(1.25)
        
        add_textbox(s5, st, step_x, step_y+0.1, step_w, 0.3, 11, C_BLUE, True, PP_ALIGN.CENTER)
        add_textbox(s5, sd, step_x, step_y+0.4, step_w, 0.6, 10, C_GRAY_TEXT, False, PP_ALIGN.CENTER)
        
        # 箭头 (原 x offset=2.6, y=4.7, w=0.5, h=0.3)
        if i < 3:
            arr_n = s5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(step_x+1.95), Inches(step_y+0.45), Inches(0.375), Inches(0.225))
            arr_n.fill.solid(); arr_n.fill.fore_color.rgb = C_GOLD
            
        step_x += step_inc

    add_script(s5, "“我们将坚持刀刃向内，持续优化模型以应对工艺变化，并通过PDCA循环，推动攻坚成果在全公司范围内的复制推广，建立长效赋能机制。”")

    # 保存
    prs.save(FILENAME)

if __name__ == "__main__":
    create_presentation()