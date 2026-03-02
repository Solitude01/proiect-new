import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement

FILENAME = "2025年党员攻坚项目汇报_AI侦测项目.pptx"

# 颜色常量
C_BLUE = RGBColor(0x00, 0x33, 0x66)
C_GOLD = RGBColor(0xB8, 0x86, 0x0B)
C_GRAY = RGBColor(0x33, 0x33, 0x33)
C_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)
C_RED = RGBColor(0xD3, 0x2F, 0x2F)
C_GREEN = RGBColor(0x28, 0xA7, 0x45)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)

def add_shadow(shape):
    """为形状添加阴影效果"""
    try:
        shadow = OxmlElement('a:effectLst')
        outer_shadow = OxmlElement('a:outerShdw')
        outer_shadow.set('blurRad', '38100')
        outer_shadow.set('dist', '38100')
        outer_shadow.set('dir', '2700000')
        outer_shadow.set('algn', 'ctr')
        
        shadow_color = OxmlElement('a:srgbClr')
        shadow_color.set('val', '000000')
        alpha = OxmlElement('a:alpha')
        alpha.set('val', '30000')
        shadow_color.append(alpha)
        
        outer_shadow.append(shadow_color)
        shadow.append(outer_shadow)
        shape._element.spPr.append(shadow)
    except:
        pass

def create_presentation():
    prs = Presentation()
    # 标准16:9尺寸
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 设置文档属性
    prs.core_properties.author = "吴佳宇"
    prs.core_properties.title = "2025年党员攻坚项目汇报_AI侦测项目"

    def add_text_box(slide, text, x, y, w, h, font_size=12, color=C_GRAY, bold=False, align=None, bg_color=None):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = shape.text_frame
        tf.word_wrap = True
        
        if bg_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg_color

        p = tf.paragraphs[0]
        if align:
            p.alignment = align
        
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.name = 'Microsoft YaHei'
        run.font.color.rgb = color
        run.font.bold = bold
        return shape

    def apply_master_layout(slide, page_type="CONTENT", slide_num=None):
        logo = slide.shapes.add_textbox(Inches(0.3), Inches(0.2 if page_type=="CONTENT" else 0.3), Inches(2), Inches(0.4))
        logo.text_frame.text = "SCC 深南电路"
        logo.text_frame.paragraphs[0].font.bold = True
        logo.text_frame.paragraphs[0].font.color.rgb = C_GOLD
        logo.text_frame.paragraphs[0].font.size = Pt(14 if page_type=="COVER" else 12)

        if page_type == "COVER":
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.225), Inches(10), Inches(0.4))
            rect.fill.solid()
            rect.fill.fore_color.rgb = C_BLUE
            rect.line.fill.background()
        else:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.75), Inches(9.4), Inches(0.015))
            line.fill.solid()
            line.fill.fore_color.rgb = C_LIGHT_GRAY
            line.line.fill.background()

            footer = slide.shapes.add_textbox(Inches(0.3), Inches(5.25), Inches(6), Inches(0.3))
            footer.text_frame.text = "2025年党员攻坚项目 | 南通深南党支部 | AI侦测辅助监控"
            footer.text_frame.paragraphs[0].font.size = Pt(8)
            footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
            
            if slide_num:
                pg = slide.shapes.add_textbox(Inches(9.4), Inches(5.25), Inches(0.3), Inches(0.3))
                pg.text_frame.text = str(slide_num)
                pg.text_frame.paragraphs[0].font.size = Pt(10)
                pg.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x88, 0x88, 0x88)
                pg.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # ==================== 幻灯片 1：封面 ====================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_layout(slide1, "COVER")

    title = add_text_box(slide1, "AI侦测辅助监控项目", 0.6, 1.8, 7.5, 1, 36, C_BLUE, True)
    subtitle = add_text_box(slide1, "2025年度党员攻坚项目汇报", 0.6, 2.7, 7.5, 0.6, 20, RGBColor(0x66,0x66,0x66))
    
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.2), Inches(1.5), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = C_GOLD
    line.line.fill.background()

    info_box = slide1.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(4.5), Inches(1.2))
    tf = info_box.text_frame
    p1 = tf.add_paragraph()
    r1 = p1.add_run()
    r1.text = "汇报支部："; r1.font.bold=True; r1.font.color.rgb=C_BLUE; r1.font.size=Pt(13)
    r2 = p1.add_run()
    r2.text = "南通深南党支部\n"; r2.font.color.rgb=C_GRAY; r2.font.size=Pt(13)
    
    p2 = tf.add_paragraph()
    r3 = p2.add_run()
    r3.text = "项目负责人："; r3.font.bold=True; r3.font.color.rgb=C_BLUE; r3.font.size=Pt(13)
    r4 = p2.add_run()
    r4.text = "吴佳宇（党员）"; r4.font.color.rgb=C_GRAY; r4.font.size=Pt(13)

    circle1 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.5), Inches(1.3), Inches(3), Inches(3))
    circle1.line.color.rgb = C_LIGHT
    circle1.line.width = Pt(20)
    circle1.fill.background()

    # ==================== 幻灯片 2：政治引领与顶层设计 ====================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_layout(slide2, "CONTENT", 1)
    
    add_text_box(slide2, "第一页：政治引领与顶层设计", 0.3, 0.4, 4.5, 0.35, 18, C_BLUE, True)
    add_text_box(slide2, "筑牢战斗堡垒，聚焦数字化瓶颈", 0.3, 0.7, 4.5, 0.3, 11, C_GRAY)

    bg_rect = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.1), Inches(9.4), Inches(3.7))
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = C_LIGHT
    bg_rect.line.fill.background()
    add_shadow(bg_rect)

    items = [
        ("🎯 政治站位", "落实上级党委战略部署，以""党建+提质增效""总方针服务高质量发展"),
        ("⚠️ 核心挑战", "传统管理模式下，人工点检有效性仅20%，制约数字化转型深化"),
        ("🔧 顶层设计", "南通深南党支部将《AI侦测辅助监控》项目定为支部级攻坚任务"),
        ("📊 关键目标", "实现点检准确率90%以上，完成AI系统对接")
    ]
    
    y_pos = 1.4
    for label, text in items:
        tb = slide2.shapes.add_textbox(Inches(0.6), Inches(y_pos), Inches(8.8), Inches(0.65))
        p = tb.text_frame.paragraphs[0]
        r1 = p.add_run(); r1.text = f"{label}："; r1.font.bold=True; r1.font.color.rgb=C_BLUE; r1.font.size=Pt(12)
        r2 = p.add_run(); r2.text = text; r2.font.color.rgb=C_GRAY; r2.font.size=Pt(11)
        y_pos += 0.85

    # ==================== 幻灯片 3：组织攻坚与党群联动 ====================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_layout(slide3, "CONTENT", 2)
    add_text_box(slide3, "第二页：组织攻坚与党群联动", 0.3, 0.4, 4.5, 0.35, 18, C_BLUE, True)
    add_text_box(slide3, "党员领题突击，群众聚力攻坚", 0.3, 0.7, 4.5, 0.3, 11, C_GRAY)

    banner = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.1), Inches(9.4), Inches(0.55))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF0)
    banner.line.color.rgb = C_RED
    add_text_box(slide3, "🚩 战斗堡垒：成立""AI数字赋能""党员攻坚突击队，实行项目化管理", 
                 0.4, 1.22, 9, 0.4, 13, C_RED, True)

    roles = [
        ("党员先锋", "吴佳宇（党员）", "项目负责人\n资源调配\n风险控制", C_BLUE),
        ("技术破题", "潘秦（党员）", "技术骨干\n预研落地\n模型有效性99%", C_BLUE),
        ("党群联动", "张旭阳（预备党员）", "现场点检\n素材采集\n数据支撑", RGBColor(0x00, 0x66, 0x99))
    ]

    x_base = 0.5
    for title, name, desc, color in roles:
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_base), Inches(1.9), Inches(2.8), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = color
        card.line.width = Pt(2)
        add_shadow(card)

        add_text_box(slide3, title, x_base, 2.05, 2.8, 0.4, 13, color, True, PP_ALIGN.CENTER)
        add_text_box(slide3, name, x_base, 2.5, 2.8, 0.4, 12, C_BLUE, True, PP_ALIGN.CENTER)
        add_text_box(slide3, desc, x_base+0.15, 2.95, 2.5, 1.1, 10, C_GRAY, False, PP_ALIGN.CENTER)

        x_base += 3.15

    add_text_box(slide3, "💡 1+N联动机制：党员带动预备党员、群众骨干（李鹏飞、丁灵），形成联合攻关合力", 
                 0.3, 4.65, 9.4, 0.4, 10, RGBColor(0x55,0x55,0x55))

    # ==================== 幻灯片 4：重点项目推进表 ====================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_layout(slide4, "CONTENT", 3)
    add_text_box(slide4, "第三页：重点项目全景推进表", 0.3, 0.4, 4.5, 0.35, 18, C_BLUE, True)
    add_text_box(slide4, "挂图作战，节点管控，实绩说话", 0.3, 0.7, 4.5, 0.3, 11, C_GRAY)

    timeline = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(8.4), Inches(0.05))
    timeline.fill.solid()
    timeline.fill.fore_color.rgb = C_LIGHT_GRAY
    timeline.line.fill.background()

    milestones = [
        ("3月15日", "现状分析", "已完成", C_GREEN),
        ("4月10日", "硬件优化", "已完成", C_GREEN),
        ("6月15日", "软件上线", "已完成", C_GREEN),
        ("9月20日", "系统测试", "已完成", C_GREEN),
        ("11月25日", "成果移交", "推进中", C_GOLD)
    ]
    
    x_positions = [1.2, 2.7, 4.2, 5.7, 7.2]
    for i, (date, title, status, color) in enumerate(milestones):
        x = x_positions[i]
        
        dot = slide4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(1.88), Inches(0.24), Inches(0.24))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.color.rgb = C_WHITE
        dot.line.width = Pt(2.5)

        add_text_box(slide4, date, x-0.3, 1.45, 0.84, 0.3, 10, color, True, PP_ALIGN.CENTER)
        add_text_box(slide4, title, x-0.3, 2.25, 0.84, 0.3, 10, C_GRAY, False, PP_ALIGN.CENTER)
        
        tag = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x-0.26), Inches(2.6), Inches(0.76), Inches(0.28))
        tag.fill.solid()
        tag.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9) if status == "已完成" else RGBColor(0xFF, 0xF3, 0xE0)
        tag.line.fill.background()
        add_text_box(slide4, status, x-0.26, 2.65, 0.76, 0.24, 9, color, True, PP_ALIGN.CENTER)

    challenge_box = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(3.3), Inches(9.4), Inches(1.1))
    challenge_box.fill.solid()
    challenge_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF9, 0xF0)
    challenge_box.line.color.rgb = C_GOLD
    add_text_box(slide4, "⚙️ 核心技术难点：模型有效性达99%（推进中）", 
                 0.5, 3.5, 8.8, 0.35, 12, C_GOLD, True)
    add_text_box(slide4, "🎯 下步重点：2025年11月25日前完成项目成果移交与规范化培训", 
                 0.5, 3.9, 8.8, 0.35, 11, C_GRAY)

    # ==================== 幻灯片 5：赋能增效与数据实绩 ====================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_layout(slide5, "CONTENT", 4)
    add_text_box(slide5, "第四页：赋能增效与数据实绩", 0.3, 0.4, 4.5, 0.35, 18, C_BLUE, True)
    add_text_box(slide5, "将党建活力转化为发展动力", 0.3, 0.7, 4.5, 0.3, 11, C_GRAY)

    box_before = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.5), Inches(2.4), Inches(2.5))
    box_before.fill.solid()
    box_before.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    box_before.line.fill.background()
    add_shadow(box_before)
    add_text_box(slide5, "改善前", 0.9, 1.65, 2.4, 0.4, 13, RGBColor(0x99,0x99,0x99), True, PP_ALIGN.CENTER)
    add_text_box(slide5, "20%", 0.9, 2.2, 2.4, 1, 44, RGBColor(0xCC,0xCC,0xCC), True, PP_ALIGN.CENTER)
    add_text_box(slide5, "人工点检准确率", 0.9, 3.5, 2.4, 0.3, 10, RGBColor(0x99,0x99,0x99), False, PP_ALIGN.CENTER)

    arrow = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.55), Inches(2.55), Inches(0.9), Inches(0.55))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = C_GOLD
    arrow.line.fill.background()

    box_after = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(1.35), Inches(2.6), Inches(2.8))
    box_after.fill.solid()
    box_after.fill.fore_color.rgb = C_WHITE
    box_after.line.color.rgb = C_BLUE
    box_after.line.width = Pt(2.5)
    add_shadow(box_after)
    add_text_box(slide5, "改善后", 4.7, 1.5, 2.6, 0.4, 13, C_BLUE, True, PP_ALIGN.CENTER)
    add_text_box(slide5, ">90%", 4.7, 2.05, 2.6, 1, 48, C_BLUE, True, PP_ALIGN.CENTER)
    add_text_box(slide5, "AI辅助监控准确率", 4.7, 3.5, 2.6, 0.3, 10, C_BLUE, False, PP_ALIGN.CENTER)

    achievements = [
        "✅ 核心指标飞跃：20%→90%+",
        "✅ 预研落地：2个AI侦测项目",
        "✅ 质量赋能：弱化异常影响",
        "✅ 长效固化：规范化流程"
    ]
    
    y = 1.5
    for text in achievements:
        add_text_box(slide5, text, 7.6, y, 2.1, 0.48, 10, C_GRAY, False)
        y += 0.58

    # ==================== 幻灯片 6：问题剖析与长效机制 ====================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_layout(slide6, "CONTENT", 5)
    add_text_box(slide6, "第五页：问题剖析与长效机制", 0.3, 0.4, 4.5, 0.35, 18, C_BLUE, True)
    add_text_box(slide6, "刀刃向内找差距，着眼长远建机制", 0.3, 0.7, 4.5, 0.3, 11, C_GRAY)

    plans = [
        ("持续优化", "应对工艺变化\n模型动态迭代\n防范风险"),
        ("成果固化", "11月25日前\n移交培训\n深度应用"),
        ("经验推广", "PDCA闭环\n党建带团建\n复制推广")
    ]

    x_base = 1.0
    for i, (title, desc) in enumerate(plans):
        oval = slide6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_base), Inches(1.6), Inches(2.2), Inches(2.2))
        oval.fill.solid()
        oval.fill.fore_color.rgb = C_WHITE
        oval.line.color.rgb = C_BLUE
        oval.line.width = Pt(2)
        add_shadow(oval)

        add_text_box(slide6, title, x_base, 2.15, 2.2, 0.4, 14, C_BLUE, True, PP_ALIGN.CENTER)
        add_text_box(slide6, desc, x_base+0.15, 2.6, 1.9, 0.9, 10, C_GRAY, False, PP_ALIGN.CENTER)

        if i < 2:
            arr = slide6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x_base+2.4), Inches(2.55), Inches(0.5), Inches(0.4))
            arr.fill.solid()
            arr.fill.fore_color.rgb = C_GOLD
            arr.line.fill.background()
        
        x_base += 3.0

    reflection_box = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(4.15), Inches(9.4), Inches(0.8))
    reflection_box.fill.solid()
    reflection_box.fill.fore_color.rgb = RGBColor(0xF0, 0xF8, 0xFF)
    reflection_box.line.color.rgb = C_BLUE
    add_text_box(slide6, "💭 反思与展望", 0.5, 4.3, 1.5, 0.3, 11, C_BLUE, True)
    add_text_box(slide6, "警惕'两张皮'思想残余，确保党建与业务深度融合；坚持PDCA闭环管理，推动成功经验在公司更广范围复制推广。", 
                 0.5, 4.6, 8.8, 0.3, 10, C_GRAY)

    # 保存文件
    prs.save(FILENAME)
    print(f"✅ PPT生成成功：{os.path.abspath(FILENAME)}")
    print(f"📊 共生成 {len(prs.slides)} 页幻灯片")
    print(f"📐 幻灯片尺寸：10×5.625英寸（标准16:9）")
    print(f"👤 作者：{prs.core_properties.author}")

if __name__ == "__main__":
    create_presentation()