import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- 全局配置 ---
FILENAME = "2025年党员攻坚项目汇报_AI侦测项目.pptx"

# 配色方案 (SCC Corporate Colors)
C_BLUE = RGBColor(0x00, 0x33, 0x66)      # 主色：深蓝
C_GOLD = RGBColor(0xB8, 0x86, 0x0B)      # 强调色：暗金
C_GRAY_TEXT = RGBColor(0x33, 0x33, 0x33) # 正文灰
C_LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)  # 浅灰背景
C_RED = RGBColor(0xD3, 0x2F, 0x2F)       # 警告/重点
C_GREEN = RGBColor(0x28, 0xA7, 0x45)     # 成功/完成
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BORDER = RGBColor(0xE0, 0xE0, 0xE0)

def create_presentation():
    # 1. 初始化 16:9 画布
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- 核心辅助函数：创建文本框 ---
    def add_textbox(slide, text, x, y, w, h, font_size=12, color=C_GRAY_TEXT, bold=False, align=None, bg_color=None):
        """快速创建一个文本框，支持背景色和对齐"""
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        
        # 背景填充
        if bg_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg_color
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        
        if align:
            p.alignment = align
            
        run = p.add_run()
        run.text = text
        run.font.name = 'Microsoft YaHei'  # 强制雅黑
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        return shape

    # --- 核心辅助函数：应用母版样式 (Logo, 页脚, 演讲稿框) ---
    def apply_master(slide, title_text, sub_text=None, page_num=None):
        # 1. 左上角 Logo
        logo = slide.shapes.add_textbox(Inches(0.4), Inches(0.4), Inches(3), Inches(0.5))
        logo.text_frame.text = "SCC深南电路"
        logo.text_frame.paragraphs[0].font.bold = True
        logo.text_frame.paragraphs[0].font.color.rgb = C_GOLD
        logo.text_frame.paragraphs[0].font.size = Pt(14)

        # 2. 标题与分割线
        add_textbox(slide, title_text, 0.4, 0.8, 8, 0.8, 24, C_BLUE, True)
        if sub_text:
             add_textbox(slide, f"|  {sub_text}", 4.5, 0.95, 6, 0.5, 16, C_GRAY_TEXT)
        
        # 分割线
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.5), Inches(12.53), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = C_GOLD
        line.line.fill.background()

        # 3. 页脚 (左侧信息 + 右侧页码)
        footer = slide.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(6), Inches(0.4))
        footer.text_frame.text = "2025年党员攻坚项目 | 南通深南党支部"
        footer.text_frame.paragraphs[0].font.size = Pt(9)
        footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

        if page_num:
            pg = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.5), Inches(0.4))
            pg.text_frame.text = str(page_num)
            pg.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            pg.text_frame.paragraphs[0].font.size = Pt(9)
            pg.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

    # --- 核心辅助函数：添加演讲稿框 (Script Box) ---
    def add_script(slide, script_text):
        # 底部灰色虚线框
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(5.8), Inches(12.53), Inches(1.1))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5) # 极浅灰
        bg.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        bg.line.dash_style = 1 # 虚线
        
        # 标签
        add_textbox(slide, "演讲关键句：", 0.5, 5.9, 2, 0.3, 11, C_GOLD, True)
        # 内容
        content = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6))
        content.text_frame.word_wrap = True
        p = content.text_frame.paragraphs[0]
        p.text = script_text
        p.font.italic = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # ==================== 封面页 (Cover) ====================
    slide_cover = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 封面背景装饰
    shape_bg = slide_cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.5), Inches(13.333), Inches(2))
    shape_bg.fill.solid()
    shape_bg.fill.fore_color.rgb = C_BLUE
    shape_bg.line.fill.background()
    
    # Logo
    logo_cover = slide_cover.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.5))
    logo_cover.text_frame.text = "SCC深南电路"
    logo_cover.text_frame.paragraphs[0].font.bold = True
    logo_cover.text_frame.paragraphs[0].font.color.rgb = C_GOLD
    logo_cover.text_frame.paragraphs[0].font.size = Pt(18)

    # 主标题
    add_textbox(slide_cover, "AI侦测辅助监控项目", 1.0, 2.5, 10, 1.5, 48, C_BLUE, True)
    add_textbox(slide_cover, "2025年度党员攻坚项目汇报", 1.0, 3.5, 8, 0.8, 24, C_GRAY_TEXT)
    
    # 金色分割线
    line_cover = slide_cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.3), Inches(1.5), Inches(0.1))
    line_cover.fill.solid()
    line_cover.fill.fore_color.rgb = C_GOLD
    line_cover.line.fill.background()

    # 底部白字信息
    add_textbox(slide_cover, "汇报支部：南通深南党支部", 1.0, 6.0, 6, 0.5, 16, C_WHITE, True)
    add_textbox(slide_cover, "项目负责人：吴佳宇", 1.0, 6.5, 6, 0.5, 16, C_WHITE, False)

    # 右侧装饰圆
    circle = slide_cover.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(1.5), Inches(5), Inches(5))
    circle.fill.background()
    circle.line.color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    circle.line.width = Pt(20)

    # ==================== P1: 项目概况 ====================
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s1, "01 项目概况", "聚焦管理瓶颈，落实数字化战略", 1)

    # 左侧文字卡片
    bg_s1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(6.5), Inches(3.8))
    bg_s1.fill.solid(); bg_s1.fill.fore_color.rgb = C_LIGHT_BG
    bg_s1.line.color.rgb = C_BORDER

    points_s1 = [
        ("政治站位", "紧密承接部门BSC战略，以数字化转型服务公司高质量发展。"),
        ("痛点分析", "传统人工点检有效性仅20%，异常发现严重滞后。"),
        ("攻坚难点", "目标准确率>90%。涉及AI模型、跨部门系统对接及现场复杂工艺。"),
        ("项目周期", "2025年03月 — 2025年11月（提质增效类）。")
    ]
    y = 2.0
    for title, desc in points_s1:
        tb = s1.shapes.add_textbox(Inches(0.7), Inches(y), Inches(6), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        r1 = p.add_run(); r1.text = f"● {title}："; r1.font.bold=True; r1.font.color.rgb=C_BLUE; r1.font.size=Pt(12)
        r2 = p.add_run(); r2.text = desc; r2.font.color.rgb=C_GRAY_TEXT; r2.font.size=Pt(12)
        y += 0.9

    # 右侧图表: 20% vs 90%
    # 标题
    add_textbox(s1, "有效性对比分析", 8.0, 1.8, 4, 0.5, 14, C_BLUE, True, PP_ALIGN.CENTER)
    
    # Bar 1 (20%)
    bar1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(4.5), Inches(1.2), Inches(1.0))
    bar1.fill.solid(); bar1.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD); bar1.line.fill.background()
    add_textbox(s1, "20%", 8.5, 4.1, 1.2, 0.4, 14, C_GRAY_TEXT, True, PP_ALIGN.CENTER)
    add_textbox(s1, "传统点检", 8.5, 5.6, 1.2, 0.4, 12, C_GRAY_TEXT, False, PP_ALIGN.CENTER)

    # Bar 2 (90%)
    bar2 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.5), Inches(2.5), Inches(1.2), Inches(3.0))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = C_BLUE; bar2.line.fill.background()
    add_textbox(s1, "90%+", 10.5, 2.1, 1.2, 0.4, 16, C_BLUE, True, PP_ALIGN.CENTER)
    add_textbox(s1, "AI辅助 (目标)", 10.5, 5.6, 1.2, 0.4, 12, C_BLUE, True, PP_ALIGN.CENTER)

    add_script(s1, "“核心目标在于突破数字化瓶颈...传统的周期性人工点检有效性仅20%...必须由党支部发挥政治引领作用，以攻坚克难的精神确保任务完成。”")

    # ==================== P2: 推进情况 ====================
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s2, "02 项目管理推进情况", "严密组织，里程碑节点按期达成", 2)

    # 组织保障 Banner
    banner_s2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(12.33), Inches(0.6))
    banner_s2.fill.solid(); banner_s2.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    banner_s2.line.color.rgb = C_BLUE
    add_textbox(s2, "🛡️ 组织保障：纳入南通深南党支部核心攻坚任务，确保资源聚焦", 0.6, 1.9, 12, 0.5, 14, C_BLUE, True)

    # 时间轴基线
    line_time = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.8), Inches(11.3), Inches(0.05))
    line_time.fill.solid(); line_time.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

    milestones = [
        ("3.15", "数据夯实", "按期完成", C_BLUE),
        ("4.10", "硬件优化", "提前完成", C_GREEN), # 高亮
        ("6.15", "软件上线", "按期上线", C_BLUE),
        ("9.20", "系统互通", "现场测试", C_BLUE)
    ]

    mx = 1.5
    for date, title, status, color in milestones:
        # 圆点
        dot = s2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(mx), Inches(3.65), Inches(0.35), Inches(0.35))
        dot.fill.solid(); dot.fill.fore_color.rgb = color
        dot.line.color.rgb = C_WHITE; dot.line.width = Pt(2)
        
        # 上方日期
        add_textbox(s2, date, mx-0.5, 3.1, 1.4, 0.4, 16, color, True, PP_ALIGN.CENTER)
        # 下方标题
        add_textbox(s2, title, mx-0.5, 4.1, 1.4, 0.4, 14, C_GRAY_TEXT, True, PP_ALIGN.CENTER)
        # 状态标签
        tag_bg = C_LIGHT_BG
        tag_font = C_GRAY_TEXT
        if status == "提前完成":
            tag_bg = RGBColor(0xE8, 0xF5, 0xE9)
            tag_font = C_GREEN
            
        tag = add_textbox(s2, status, mx-0.6, 4.6, 1.6, 0.4, 10, tag_font, True, PP_ALIGN.CENTER, tag_bg)
        mx += 3.0

    add_script(s2, "“团队在4月10日即完成硬件优化，比计划提前了近半个月...其中软件开发级上线应用已于6月15日高质量完成...展现了项目组令行禁止的风貌。”")

    # ==================== P3: 党员作用 ====================
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s3, "03 党员攻坚作用", "战斗堡垒领航，先锋模范先行", 3)

    # 红色突击队 Banner
    banner_red = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(12.33), Inches(0.6))
    banner_red.fill.solid(); banner_red.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
    banner_red.line.color.rgb = C_RED
    add_textbox(s3, "🚩 支部行动：成立“AI数字赋能”党员突击队，高位协调，解决跨部门协同难题", 0.6, 1.9, 12, 0.5, 14, C_RED, True)

    # 3张人物卡片
    roles = [
        ("项目负责人", "吴佳宇 (党员)", "负责资源调配与风险控制。\n在6月关键节点果断决策调整预研方向。"),
        ("技术骨干", "潘秦 (党员)", "主导AI模型实现，攻克技术壁垒。\n克服“素材采集困难”风险，确保指标达成。"),
        ("一线先锋", "张旭阳 (预备党员)", "冲锋一线，配合完成现场点检及大量素材采集，\n夯实数据基础。")
    ]

    cx = 0.5
    for r_title, r_name, r_desc in roles:
        # 卡片头
        head = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(2.8), Inches(3.9), Inches(0.6))
        head.fill.solid(); head.fill.fore_color.rgb = C_BLUE; head.line.fill.background()
        add_textbox(s3, r_title, cx, 2.9, 3.9, 0.5, 12, RGBColor(0xAD, 0xD8, 0xE6), True, PP_ALIGN.CENTER)
        
        # 卡片身
        body = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(3.4), Inches(3.9), Inches(2.0))
        body.fill.solid(); body.fill.fore_color.rgb = C_LIGHT_BG; body.line.fill.background()
        
        add_textbox(s3, r_name, cx, 3.6, 3.9, 0.5, 14, C_BLUE, True, PP_ALIGN.CENTER)
        add_textbox(s3, r_desc, cx+0.1, 4.1, 3.7, 1.2, 11, C_GRAY_TEXT, False, PP_ALIGN.CENTER)
        cx += 4.2

    add_script(s3, "“吴佳宇同志果断决策...技术骨干潘秦同志勇挑重担...预备党员张旭阳同志冲锋一线...正是有了党员带头冲锋，才得以攻克技术难关。”")

    # ==================== P4: 效益成果 ====================
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s4, "04 项目效益和成果", "质量跃升，长效机制固化", 4)

    # 左侧对比框
    # Before
    box_b = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.5), Inches(3.0), Inches(2.0))
    box_b.fill.solid(); box_b.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE); box_b.line.fill.background()
    add_textbox(s4, "20%", 1.0, 2.8, 3.0, 0.8, 36, RGBColor(0x99, 0x99, 0x99), True, PP_ALIGN.CENTER)
    add_textbox(s4, "Before (传统点检)", 1.0, 3.8, 3.0, 0.5, 12, RGBColor(0x66, 0x66, 0x66), False, PP_ALIGN.CENTER)

    # 箭头
    arr = s4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(3.2), Inches(1.0), Inches(0.6))
    arr.fill.solid(); arr.fill.fore_color.rgb = C_GOLD; arr.line.fill.background()

    # After
    box_a = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(2.2), Inches(3.5), Inches(2.6))
    box_a.fill.solid(); box_a.fill.fore_color.rgb = C_WHITE
    box_a.line.color.rgb = C_BLUE; box_a.line.width = Pt(3)
    add_textbox(s4, ">90%", 5.5, 2.8, 3.5, 0.8, 44, C_BLUE, True, PP_ALIGN.CENTER)
    add_textbox(s4, "After (AI辅助)", 5.5, 4.0, 3.5, 0.5, 14, C_BLUE, True, PP_ALIGN.CENTER)

    # 右侧价值列表
    v_bg = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(2.2), Inches(3.3), Inches(2.6))
    v_bg.fill.solid(); v_bg.fill.fore_color.rgb = C_LIGHT_BG; v_bg.line.fill.background()
    
    add_textbox(s4, "价值与成果", 9.5, 2.4, 3.3, 0.5, 14, C_BLUE, True, PP_ALIGN.CENTER)
    vals = "✅ 全面对接AI侦测系统\n✅ 落地2个辅助监控项目\n✅ 解放人力，弱化人为影响\n✅ 机制固化，纳入标准流程"
    add_textbox(s4, vals, 9.7, 3.0, 2.9, 1.5, 11, C_GRAY_TEXT)

    add_script(s4, "“本项目实现了质量管理模式的代际升级...彻底解决了异常发现滞后的问题...成功将一次攻坚行动转化为一套长效运行的制度体系。”")

    # ==================== P5: 后续计划 ====================
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s5, "05 其他需说明事项", "巩固成果，持续深化数字赋能", 5)

    steps = [
        ("持续迭代", "针对工艺变化\n优化模型有效性"),
        ("成果移交", "11月25日前\n完成培训与移交"),
        ("复制推广", "党建+技术融合\n推广至其他板块")
    ]

    sx = 1.5
    for i, (title, desc) in enumerate(steps):
        # 圆环
        oval = s5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(sx), Inches(2.2), Inches(2.5), Inches(2.5))
        oval.fill.solid(); oval.fill.fore_color.rgb = C_WHITE
        oval.line.color.rgb = C_BLUE; oval.line.width = Pt(3)
        
        add_textbox(s5, title, sx, 2.8, 2.5, 0.5, 16, C_BLUE, True, PP_ALIGN.CENTER)
        add_textbox(s5, desc, sx, 3.4, 2.5, 1.0, 12, C_GRAY_TEXT, False, PP_ALIGN.CENTER)

        # 箭头
        if i < 2:
            arr_s = s5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(sx+2.6), Inches(3.3), Inches(0.8), Inches(0.4))
            arr_s.fill.solid(); arr_s.fill.fore_color.rgb = C_GOLD; arr_s.line.fill.background()

        sx += 3.8

    add_script(s5, "“数字化没有终点...AI模型仍需持续迭代...我们致力于将本次攻坚经验推广应用，为公司高质量发展贡献持续的红色力量！”")

    # 保存文件
    try:
        prs.save(FILENAME)
        print(f"✅ 成功生成文件：{os.path.abspath(FILENAME)}")
        print("💡 提示：请直接在当前目录查找 pptx 文件。")
    except Exception as e:
        print(f"❌ 生成失败：{e}")

if __name__ == "__main__":
    create_presentation()