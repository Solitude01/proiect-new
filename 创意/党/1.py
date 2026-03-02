import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR

# --- 配置部分 ---
FILENAME = "2025年党员攻坚项目汇报_AI侦测项目.pptx"

# 颜色常量 (RGB)
C_BLUE = RGBColor(0x00, 0x33, 0x66)
C_GOLD = RGBColor(0xB8, 0x86, 0x0B)
C_GRAY = RGBColor(0x33, 0x33, 0x33)
C_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)
C_RED = RGBColor(0xD3, 0x2F, 0x2F)
C_GREEN = RGBColor(0x28, 0xA7, 0x45)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)

def create_presentation():
    # 1. 初始化 PPT，设置为 16:9
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- 辅助函数：添加文本框 ---
    def add_text_box(slide, text, x, y, w, h, font_size=12, color=C_GRAY, bold=False, align=None, bg_color=None):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = shape.text_frame
        tf.word_wrap = True
        
        # 如果有背景色
        if bg_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg_color

        p = tf.paragraphs[0]
        if align:
            p.alignment = align
        
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.name = 'Microsoft YaHei' # 尝试设置字体，系统需安装
        run.font.color.rgb = color
        run.font.bold = bold
        return shape

    # --- 辅助函数：绘制母版元素 (Logo, 分割线, 页脚) ---
    def apply_master_layout(slide, page_type="CONTENT", slide_num=None):
        # 左上角 Logo
        logo = slide.shapes.add_textbox(Inches(0.4), Inches(0.3 if page_type=="CONTENT" else 0.4), Inches(3), Inches(0.5))
        logo.text_frame.text = "SCC深南电路"
        logo.text_frame.paragraphs[0].font.bold = True
        logo.text_frame.paragraphs[0].font.color.rgb = C_GOLD
        logo.text_frame.paragraphs[0].font.size = Pt(16 if page_type=="COVER" else 14)

        if page_type == "COVER":
            # 封面底部蓝条
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5))
            rect.fill.solid()
            rect.fill.fore_color.rgb = C_BLUE
            rect.line.fill.background() # 无边框
        else:
            # 正文页分割线
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.0), Inches(12.53), Inches(0.02))
            line.fill.solid()
            line.fill.fore_color.rgb = C_LIGHT_GRAY
            line.line.fill.background()

            # 页脚
            footer = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(5), Inches(0.5))
            footer.text_frame.text = "2025年党员攻坚项目 | 南通深南党支部"
            footer.text_frame.paragraphs[0].font.size = Pt(10)
            footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
            
            # 页码 (模拟)
            if slide_num:
                pg = slide.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.5), Inches(0.5))
                pg.text_frame.text = str(slide_num)
                pg.text_frame.paragraphs[0].font.size = Pt(10)
                pg.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
                pg.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # --- 辅助函数：添加演讲脚本框 ---
    def add_script_box(slide, text):
        # 虚线框背景
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(5.8), Inches(12.53), Inches(1.0))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
        rect.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        rect.line.dash_style = 1 # 简单的点线

        # 标签
        lbl = slide.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(2), Inches(0.3))
        lbl.text_frame.text = "演讲关键句："
        lbl.text_frame.paragraphs[0].font.bold = True
        lbl.text_frame.paragraphs[0].font.color.rgb = C_GOLD
        lbl.text_frame.paragraphs[0].font.size = Pt(12)

        # 内容
        content = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5))
        content.text_frame.word_wrap = True
        p = content.text_frame.paragraphs[0]
        p.text = text
        p.font.italic = True
        p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        p.font.size = Pt(12)

    # ==================== 幻灯片 1：封面 ====================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6]) # 6 is Blank
    apply_master_layout(slide1, "COVER")

    # 标题
    title = add_text_box(slide1, "AI侦测辅助监控项目", 0.8, 2.5, 10, 1.5, 44, C_BLUE, True)
    subtitle = add_text_box(slide1, "2025年度党员攻坚项目汇报", 0.8, 3.5, 10, 1, 24, RGBColor(0x66,0x66,0x66))
    
    # 金色短线
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.2), Inches(1.5), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = C_GOLD
    line.line.fill.background()

    # 汇报信息
    info_box = slide1.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(6), Inches(2))
    tf = info_box.text_frame
    p1 = tf.add_paragraph()
    r1 = p1.add_run()
    r1.text = "汇报支部："; r1.font.bold=True; r1.font.color.rgb=C_BLUE; r1.font.size=Pt(16)
    r2 = p1.add_run()
    r2.text = "南通深南党支部\n"; r2.font.color.rgb=C_GRAY; r2.font.size=Pt(16)
    
    p2 = tf.add_paragraph()
    r3 = p2.add_run()
    r3.text = "项目负责人："; r3.font.bold=True; r3.font.color.rgb=C_BLUE; r3.font.size=Pt(16)
    r4 = p2.add_run()
    r4.text = "吴佳宇"; r4.font.color.rgb=C_GRAY; r4.font.size=Pt(16)

    # 装饰圆环
    circle1 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(1.5), Inches(4.5), Inches(4.5))
    circle1.line.color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    circle1.line.width = Pt(30)
    circle1.fill.background()

    # ==================== 幻灯片 2：项目概况 ====================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_layout(slide2, "CONTENT", 1)
    
    # 页面标题
    add_text_box(slide2, "01 项目概况", 0.4, 0.5, 4, 0.8, 20, C_BLUE, True)
    add_text_box(slide2, "|  聚焦管理瓶颈，落实数字化战略", 2.5, 0.6, 6, 0.8, 14, C_GRAY)

    # 左侧灰色背景块
    bg_rect = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.5), Inches(7.5), Inches(4.0))
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = C_LIGHT
    bg_rect.line.fill.background()

    # 列表内容
    items = [
        ("政治站位", "紧密承接部门BSC战略，以数字化转型服务公司高质量发展。"),
        ("痛点分析", "传统人工点检有效性仅20%，异常发现严重滞后。"),
        ("攻坚难点", "目标准确率>90%。涉及AI模型、跨部门系统对接及现场复杂工艺。"),
        ("项目周期", "2025年03月 — 2025年11月（提质增效类）。")
    ]
    y_pos = 1.7
    for label, text in items:
        tb = slide2.shapes.add_textbox(Inches(0.6), Inches(y_pos), Inches(7), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        r1 = p.add_run(); r1.text = f"● {label}："; r1.font.bold=True; r1.font.color.rgb=C_BLUE; r1.font.size=Pt(14)
        r2 = p.add_run(); r2.text = text; r2.font.color.rgb=C_GRAY; r2.font.size=Pt(14)
        y_pos += 0.9

    # 右侧对比图
    # 20%
    bar1 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(3.5), Inches(0.8), Inches(0.8))
    bar1.fill.solid(); bar1.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC); bar1.line.fill.background()
    add_text_box(slide2, "20%", 8.5, 3.1, 0.8, 0.5, 14, C_GRAY, True, PP_ALIGN.CENTER)
    add_text_box(slide2, "传统点检", 8.4, 4.4, 1, 0.5, 12, C_GRAY, False, PP_ALIGN.CENTER)

    # 90%
    bar2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.5), Inches(2.5), Inches(0.8), Inches(1.8))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = C_BLUE; bar2.line.fill.background()
    add_text_box(slide2, "90%+", 10.5, 2.1, 0.8, 0.5, 14, C_BLUE, True, PP_ALIGN.CENTER)
    add_text_box(slide2, "AI辅助", 10.4, 4.4, 1, 0.5, 12, C_BLUE, False, PP_ALIGN.CENTER)

    add_script_box(slide2, "“传统的周期性人工点检有效性仅20%...必须由党支部发挥政治引领作用，以攻坚克难的精神确保任务完成。”")

    # ==================== 幻灯片 3：推进情况 ====================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_layout(slide3, "CONTENT", 2)
    add_text_box(slide3, "02 项目管理推进情况", 0.4, 0.5, 5, 0.8, 20, C_BLUE, True)

    add_text_box(slide3, "🛡️ 组织保障：纳入南通深南党支部核心攻坚任务，确保资源聚焦", 0.4, 1.3, 12, 0.5, 16, C_BLUE, True)

    # 时间轴线
    timeline = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.5), Inches(11), Inches(0.05))
    timeline.fill.solid(); timeline.fill.fore_color.rgb = RGBColor(0xDD,0xDD,0xDD)

    milestones = [
        ("3.15", "数据夯实", "按期", C_BLUE),
        ("4.10", "硬件优化", "提前完成", C_GREEN),
        ("6.15", "软件上线", "按期", C_BLUE),
        ("9.20", "系统互通", "按期", C_BLUE)
    ]
    
    x_base = 1.5
    for date, title, status, color in milestones:
        # 圆点
        dot = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_base), Inches(3.35), Inches(0.35), Inches(0.35))
        dot.fill.solid(); dot.fill.fore_color.rgb = color
        dot.line.color.rgb = C_WHITE; dot.line.width = Pt(2)

        # 文字
        add_text_box(slide3, date, x_base-0.5, 2.8, 1.4, 0.5, 16, color, True, PP_ALIGN.CENTER)
        add_text_box(slide3, title, x_base-0.5, 3.8, 1.4, 0.5, 14, C_GRAY, False, PP_ALIGN.CENTER)
        
        # 状态标签
        if status == "提前完成":
            tag = add_text_box(slide3, status, x_base-0.6, 4.3, 1.6, 0.4, 11, C_GREEN, True, PP_ALIGN.CENTER)
            tag.fill.solid(); tag.fill.fore_color.rgb = RGBColor(0xE8,0xF5,0xE9)

        x_base += 3.0

    add_script_box(slide3, "“团队在4月10日即完成硬件优化，比计划提前了近半个月...展现了项目组令行禁止的攻坚风貌。”")

    # ==================== 幻灯片 4：党员作用 ====================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_layout(slide4, "CONTENT", 3)
    add_text_box(slide4, "03 党员攻坚作用", 0.4, 0.5, 5, 0.8, 20, C_BLUE, True)

    # 红色Banner
    banner = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.8))
    banner.fill.solid(); banner.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF0)
    banner.line.color.rgb = C_RED
    add_text_box(slide4, "🚩 支部行动：成立“AI数字赋能”党员突击队，高位协调，解决跨部门协同难题", 
                 0.5, 1.45, 12, 0.6, 16, C_RED, True)

    roles = [
        ("项目负责人", "吴佳宇 (党员)", "统筹资源 / 风险控制\n在6月关键节点果断调整方向"),
        ("技术骨干", "潘秦 (党员)", "主导模型 / 攻克壁垒\n克服素材采集困难，确保指标达成"),
        ("一线先锋", "张旭阳 (预备党员)", "现场点检 / 数据采集\n冲锋一线，夯实数据基础")
    ]

    x_base = 0.5
    for role_title, name, desc in roles:
        # 标题栏
        header = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_base), Inches(2.5), Inches(3.8), Inches(0.6))
        header.fill.solid(); header.fill.fore_color.rgb = C_BLUE
        add_text_box(slide4, role_title, x_base, 2.6, 3.8, 0.6, 14, RGBColor(173,216,230), False, PP_ALIGN.CENTER)

        # 内容框
        box = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_base), Inches(3.1), Inches(3.8), Inches(2.0))
        box.fill.solid(); box.fill.fore_color.rgb = C_LIGHT
        box.line.fill.background() # 无边框

        add_text_box(slide4, name, x_base, 3.3, 3.8, 0.5, 16, C_BLUE, True, PP_ALIGN.CENTER)
        add_text_box(slide4, desc, x_base+0.1, 3.8, 3.6, 1.2, 12, C_GRAY, False, PP_ALIGN.CENTER)

        x_base += 4.1

    add_script_box(slide4, "“正是有了党员的带头冲锋、主动担当，我们才得以攻克技术难关...实现预定目标。”")

    # ==================== 幻灯片 5：效益成果 ====================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_layout(slide5, "CONTENT", 4)
    add_text_box(slide5, "04 项目效益和成果", 0.4, 0.5, 5, 0.8, 20, C_BLUE, True)

    # Before
    box_b = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.0), Inches(3), Inches(2.5))
    box_b.fill.solid(); box_b.fill.fore_color.rgb = RGBColor(0xEE,0xEE,0xEE)
    box_b.line.fill.background()
    add_text_box(slide5, "20%", 1.5, 2.5, 3, 1, 40, RGBColor(0x99,0x99,0x99), True, PP_ALIGN.CENTER)
    add_text_box(slide5, "Before", 1.5, 3.5, 3, 0.5, 14, RGBColor(0x99,0x99,0x99), False, PP_ALIGN.CENTER)

    # Arrow
    arrow = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.8), Inches(3.0), Inches(1), Inches(0.6))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = C_GOLD; arrow.line.fill.background()

    # After
    box_a = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.2), Inches(1.8), Inches(3.5), Inches(3.0))
    box_a.fill.solid(); box_a.fill.fore_color.rgb = C_WHITE
    box_a.line.color.rgb = C_BLUE; box_a.line.width = Pt(3)
    add_text_box(slide5, ">90%", 6.2, 2.5, 3.5, 1, 48, C_BLUE, True, PP_ALIGN.CENTER)
    add_text_box(slide5, "After (质的飞跃)", 6.2, 3.8, 3.5, 0.5, 14, C_BLUE, False, PP_ALIGN.CENTER)

    # 价值列表
    val_box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.2), Inches(2.0), Inches(2.5), Inches(2.5))
    val_box.fill.solid(); val_box.fill.fore_color.rgb = C_LIGHT
    add_text_box(slide5, "价值维度", 10.2, 2.2, 2.5, 0.5, 14, C_BLUE, True, PP_ALIGN.CENTER)
    vals = "• 全面对接\n• 落地2项\n• 解放人力\n• 机制固化"
    add_text_box(slide5, vals, 10.4, 2.8, 2.3, 1.5, 12, C_GRAY)

    add_script_box(slide5, "“本项目最大的成果是实现了质量管理模式的代际升级...成功将一次攻坚行动转化为一套长效运行的制度体系。”")

    # ==================== 幻灯片 6：其他事项 ====================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_layout(slide6, "CONTENT", 5)
    add_text_box(slide6, "05 其他需说明事项", 0.4, 0.5, 5, 0.8, 20, C_BLUE, True)

    plans = [
        ("持续迭代", "针对工艺变化\n优化模型有效性"),
        ("成果移交", "11月25日前\n完成移交与培训"),
        ("复制推广", "党建+技术融合\n推广至其他板块")
    ]

    x_base = 1.5
    for i, (title, desc) in enumerate(plans):
        # 圆圈
        oval = slide6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_base), Inches(2.0), Inches(2.5), Inches(2.5))
        oval.fill.solid(); oval.fill.fore_color.rgb = C_WHITE
        oval.line.color.rgb = C_BLUE; oval.line.width = Pt(2)

        add_text_box(slide6, title, x_base, 2.6, 2.5, 0.5, 18, C_BLUE, True, PP_ALIGN.CENTER)
        add_text_box(slide6, desc, x_base, 3.2, 2.5, 1.0, 14, C_GRAY, False, PP_ALIGN.CENTER)

        # 箭头
        if i < 2:
            arrow = slide6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x_base+2.8), Inches(3.0), Inches(0.5), Inches(0.4))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = C_GOLD; arrow.line.fill.background()
        
        x_base += 3.8

    add_script_box(slide6, "“数字化没有终点...我们致力于将本次攻坚经验推广应用，为公司高质量发展贡献持续的红色力量！”")

    # 保存
    try:
        prs.save(FILENAME)
        print(f"✅ 成功生成文件：{os.path.abspath(FILENAME)}")
    except Exception as e:
        print(f"❌ 生成失败：{e}")

if __name__ == "__main__":
    create_presentation()