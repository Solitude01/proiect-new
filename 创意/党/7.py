import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- 全局配置 ---
# 文件名
FILENAME = "2025年党员攻坚项目汇报_AI侦测项目7.pptx"
AUTHOR = "吴佳宇"

# 配色方案 (SCC Corporate Identity)
C_BLUE = RGBColor(0, 51, 102)        # #003366 主色 (SCC深蓝)
C_GOLD = RGBColor(184, 134, 11)      # #B8860B 辅助色 (暗金)
C_GRAY_TEXT = RGBColor(51, 51, 51)   # #333333 正文 (深灰)
C_LIGHT_BG = RGBColor(248, 249, 250) # #F8F9FA 卡片背景
C_RED = RGBColor(211, 47, 47)        # #D32F2F 警告/重点
C_GREEN = RGBColor(40, 167, 69)      # #28A745 完成/正面
C_WHITE = RGBColor(255, 255, 255)    # #FFFFFF
C_BORDER = RGBColor(224, 224, 224)   # #E0E0E0

def create_presentation():
    # 1. 初始化 PPT画布
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 设置核心属性
    prs.core_properties.author = AUTHOR
    prs.core_properties.title = "AI侦测辅助监控项目汇报"

    # --- 核心辅助函数 ---
    def add_textbox(slide, text, x, y, w, h, font_size=12, color=C_GRAY_TEXT, bold=False, align=None, bg_color=None, border_color=None):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        if bg_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if align:
            p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = 'Microsoft YaHei'
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        return shape

    def add_card(slide, title, content, x, y, w, h, highlight=False):
        # 绘制卡片背景
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_LIGHT_BG
        shape.line.color.rgb = C_GOLD if highlight else C_BLUE
        shape.line.width = Pt(1.5 if highlight else 0.75)
        
        # 添加阴影效果
        shadow = shape.shadow
        shadow.inherit = False
        shadow.blur_radius = Pt(3)
        shadow.distance = Pt(2)
        shadow.angle = 45
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+0.1), Inches(w-0.2), Inches(0.4))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = title
        tp.font.bold = True
        tp.font.size = Pt(12)
        tp.font.name = 'Microsoft YaHei'
        tp.font.color.rgb = C_BLUE
        
        # 内容
        content_box = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+0.45), Inches(w-0.2), Inches(h-0.55))
        content_box.text_frame.word_wrap = True
        cp = content_box.text_frame.paragraphs[0]
        cp.text = content
        cp.font.size = Pt(10)
        cp.font.name = 'Microsoft YaHei'
        cp.font.color.rgb = C_GRAY_TEXT

    # --- 母版应用函数 ---
    def apply_master(slide, title_text, sub_text=None, page_num=None):
        # 1. Logo
        logo = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(3), Inches(0.5))
        logo.text_frame.text = "SCC深南电路"
        logo.text_frame.paragraphs[0].font.bold = True
        logo.text_frame.paragraphs[0].font.color.rgb = C_GOLD
        logo.text_frame.paragraphs[0].font.size = Pt(14)

        # 2. 标题与分割线
        add_textbox(slide, title_text, 0.3, 0.5, 8.5, 0.6, 20, C_BLUE, True)
        if sub_text:
             add_textbox(slide, f"|  {sub_text}", 0.3, 0.9, 8, 0.3, 11, C_GRAY_TEXT)
        
        # 分割线
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.2), Inches(9.4), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = C_GOLD
        line.line.fill.background()

        # 3. 页脚
        footer_y = 5.35
        footer = slide.shapes.add_textbox(Inches(0.3), Inches(footer_y), Inches(5), Inches(0.3))
        footer.text_frame.text = "2025年党员攻坚项目 | 南通深南党支部"
        footer.text_frame.paragraphs[0].font.size = Pt(9)
        footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(170, 170, 170)

        if page_num:
            pg = slide.shapes.add_textbox(Inches(9.2), Inches(footer_y), Inches(0.5), Inches(0.3))
            pg.text_frame.text = str(page_num)
            pg.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            pg.text_frame.paragraphs[0].font.size = Pt(9)
            pg.text_frame.paragraphs[0].font.color.rgb = RGBColor(170, 170, 170)

    # --- 演讲稿框函数 ---
    def add_script(slide, script_text):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(4.3), Inches(9.4), Inches(1.0))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(245, 245, 245)
        bg.line.color.rgb = RGBColor(204, 204, 204)
        bg.line.dash_style = 1 
        
        add_textbox(slide, "📢 演讲逐字稿摘要：", 0.4, 4.35, 2, 0.3, 10, C_GOLD, True)
        content = slide.shapes.add_textbox(Inches(0.4), Inches(4.55), Inches(9.2), Inches(0.7))
        content.text_frame.word_wrap = True
        p = content.text_frame.paragraphs[0]
        p.text = script_text
        p.font.italic = True
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(80, 80, 80)

    # ==================== 封面页 (Cover) ====================
    slide_cover = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景色块
    shape_bg = slide_cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.125), Inches(10), Inches(1.5))
    shape_bg.fill.solid(); shape_bg.fill.fore_color.rgb = C_BLUE; shape_bg.line.fill.background()
    
    # Logo
    logo_cover = slide_cover.shapes.add_textbox(Inches(0.375), Inches(0.375), Inches(3), Inches(0.5))
    logo_cover.text_frame.text = "SCC深南电路"
    logo_cover.text_frame.paragraphs[0].font.bold = True
    logo_cover.text_frame.paragraphs[0].font.color.rgb = C_GOLD
    logo_cover.text_frame.paragraphs[0].font.size = Pt(18)

    # 标题
    add_textbox(slide_cover, "AI侦测辅助监控", 0.75, 1.875, 8, 1.2, 40, C_BLUE, True)
    add_textbox(slide_cover, "项目类型：提质增效类 | 周期：2025年03月—2025年11月", 0.75, 2.7, 8, 0.6, 16, C_GRAY_TEXT)
    
    # 分割线
    line_cover = slide_cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(3.3), Inches(1.125), Inches(0.08))
    line_cover.fill.solid(); line_cover.fill.fore_color.rgb = C_GOLD; line_cover.line.fill.background()

    # 底部信息
    add_textbox(slide_cover, "汇报支部：南通深南党支部", 0.75, 4.5, 5, 0.4, 14, C_WHITE, True)
    add_textbox(slide_cover, "项目负责人：吴佳宇", 0.75, 4.9, 5, 0.4, 14, C_WHITE, False)

    # 装饰圆环
    circle = slide_cover.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.75), Inches(1.125), Inches(3.75), Inches(3.75))
    circle.fill.background(); circle.line.color.rgb = RGBColor(240, 240, 240); circle.line.width = Pt(15)

    # ==================== 第 1 页：项目概况 ====================
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s1, "项目概况：政治引领与顶层设计", "聚焦数字化瓶颈，落实“党建+数字赋能”", 1)

    card_y = 1.4
    card_w = 4.5
    card_h = 1.4
    
    # 左卡片：政治站位
    add_card(s1, "🚩 政治站位与方针", 
             "● 落实部门BSC战略，以“党建+数字赋能”为总方针。\n● 旨在将党组织核心作用体现在解决生产经营核心难题上。\n● 服务高质量发展，实现提质增效。", 
             0.375, card_y, card_w, card_h)
    
    # 右卡片：核心挑战
    add_card(s1, "⚠️ 核心挑战 (痛点)", 
             "● 传统人工点检有效性仅占20%。\n● 异常发现滞后，严重制约现场质量改善。\n● 核心业务流程中的突出“卡脖子”瓶颈。", 
             5.1, card_y, card_w, card_h, highlight=True)

    # 底部目标区域
    goal_bg_y = 3.0
    goal_bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.375), Inches(goal_bg_y), Inches(9.225), Inches(1.1))
    goal_bg.fill.solid(); goal_bg.fill.fore_color.rgb = C_LIGHT_BG
    goal_bg.line.color.rgb = C_BLUE
    
    add_textbox(s1, "🎯 攻坚目标", 0.45, 3.1, 4, 0.3, 12, C_BLUE, True)
    
    # 目标 1
    add_textbox(s1, "90%+", 0.75, 3.4, 1.5, 0.6, 26, C_BLUE, True, PP_ALIGN.CENTER)
    add_textbox(s1, "点检准确率", 0.75, 3.8, 1.5, 0.3, 10, C_GRAY_TEXT, False, PP_ALIGN.CENTER)
    
    # 目标 2
    add_textbox(s1, "2个", 3.0, 3.4, 1.5, 0.6, 26, C_BLUE, True, PP_ALIGN.CENTER)
    add_textbox(s1, "预研AI项目", 3.0, 3.8, 1.5, 0.3, 10, C_GRAY_TEXT, False, PP_ALIGN.CENTER)

    # 描述
    add_textbox(s1, "通过引入AI侦测，完成辅助监控系统对接，挑战将点检准确率提升至90%以上的硬性指标。", 
                5.0, 3.4, 4.4, 0.6, 11, C_GRAY_TEXT)

    s1_script = ('"尊敬的各位领导、各位评委：本次汇报的项目《AI侦测辅助监控》，是南通深南党支部坚决贯彻落实公司数字化战略部署，实现提质增效的关键实践。\n'
                 '项目立项之初，我们直面核心业务流程中传统周期性点检有效性仅20%的痛点。这是一个必须由党员带头突破的‘卡脖子’瓶颈。\n'
                 '我们确立了攻坚目标：不仅要完成系统全面对接，更要挑战将点检准确率提升至90%以上，确保党组织成为突破数字化瓶颈的战斗堡垒。"').replace('\n', ' ')
    add_script(s1, s1_script)

    # ==================== 第 2 页：项目管理推进情况 ====================
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s2, "项目管理推进情况：严密组织", "里程碑节点按期达成，挂图作战", 2)

    # 左侧：组织保障
    add_textbox(s2, "🏢 组织保障", 0.375, 1.35, 3, 0.3, 12, C_BLUE, True)
    org_box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.375), Inches(1.7), Inches(2.5), Inches(1.5))
    org_box.fill.solid(); org_box.fill.fore_color.rgb = C_LIGHT_BG; org_box.line.color.rgb = C_BLUE
    add_textbox(s2, "项目负责人：吴佳宇\n\n项目性质：\n支部级核心任务\n\n管理机制：\n挂图作战，节点精准", 
                0.5, 1.8, 2.3, 1.3, 10, C_GRAY_TEXT)

    # 右侧：时间轴列表
    add_textbox(s2, "📅 关键时间节点达成 (实绩)", 3.1, 1.35, 4, 0.3, 12, C_BLUE, True)
    
    # 列表头
    headers = ["时间节点", "里程碑内容", "状态"]
    cols_x = [3.1, 5.1, 8.5]
    cols_w = [1.8, 3.2, 1.2]
    
    header_y = 1.7
    for i, h in enumerate(headers):
        s2.shapes.add_textbox(Inches(cols_x[i]), Inches(header_y), Inches(cols_w[i]), Inches(0.3)).text_frame.text = h
        
    # 数据行
    milestones = [
        ("23/03/15", "现状数据确认及分析", "🚀 提前完成"),
        ("23/04/10", "硬件调整及优化\n现状问题点改善", "🚀 提前完成"),
        ("23/06/15", "辅助系统软件开发级\n上线应用", "✅ 按期完成"),
        ("23/09/20", "互通模型及系统应用\n完成现场测试", "✅ 按期完成"),
        ("当前状态", "核心环节“模型实现”\n与“移交现场管理”", "⏳ 正在进行")
    ]
    
    y_pos = 2.1
    for date, task, status in milestones:
        # 背景条
        bg_bar = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.1), Inches(y_pos), Inches(6.6), Inches(0.45))
        bg_bar.fill.solid(); bg_bar.fill.fore_color.rgb = C_WHITE
        bg_bar.line.color.rgb = C_BORDER
        
        # 内容
        add_textbox(s2, date, cols_x[0]+0.1, y_pos+0.05, cols_w[0], 0.4, 10, C_GRAY_TEXT, True)
        add_textbox(s2, task, cols_x[1]+0.1, y_pos+0.05, cols_w[1], 0.4, 9, C_GRAY_TEXT)
        
        # 状态颜色
        st_color = C_GREEN if "完成" in status else (C_BLUE if "进行" in status else C_GRAY_TEXT)
        add_textbox(s2, status, cols_x[2]+0.1, y_pos+0.05, cols_w[2], 0.4, 9, st_color, True)
        
        y_pos += 0.52

    s2_script = ('"本项目在南通深南党支部的统一领导下，坚持严密的组织体系和挂图作战机制。\n'
                 '请看我们的进度全景：数据确认（3月15日）和硬件调整（4月10日）均提前完成。辅助系统软件6月15日按期上线。互通模型现场测试9月20日成功验证。\n'
                 '目前，正全力冲刺模型实现和成果移交两大核心环节，所有核心里程碑均高质量达成。"').replace('\n', ' ')
    add_script(s2, s2_script)

    # ==================== 第 3 页：党员攻坚作用 ====================
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s3, "党员攻坚作用：战斗堡垒领航", "先锋模范先行，党群联动攻关", 3)

    # 顶部：支部堡垒
    fortress_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(1.3), Inches(6), Inches(0.6))
    fortress_box.fill.solid(); fortress_box.fill.fore_color.rgb = C_BLUE
    add_textbox(s3, "🛡️ 支部堡垒：成立“AI数字赋能”党员突击队", 2, 1.4, 6, 0.4, 12, C_WHITE, True, PP_ALIGN.CENTER)
    add_textbox(s3, "高位协调，解决跨部门协同难题", 2, 1.7, 6, 0.3, 9, C_WHITE, False, PP_ALIGN.CENTER)

    # 连接线
    connector = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.95), Inches(1.9), Inches(0.1), Inches(0.3))
    connector.fill.solid(); connector.fill.fore_color.rgb = C_GOLD

    # 下方三列布局
    box_y = 2.2
    box_w = 2.9
    box_h = 2.0
    gap = 0.3
    start_x = 0.5

    # 1. 党员先锋
    box1 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(start_x), Inches(box_y), Inches(box_w), Inches(box_h))
    box1.fill.solid(); box1.fill.fore_color.rgb = C_LIGHT_BG; box1.line.color.rgb = C_BLUE
    add_textbox(s3, "👤 党员先锋：吴佳宇", start_x+0.1, box_y+0.1, box_w-0.2, 0.3, 11, C_BLUE, True)
    add_textbox(s3, "负责整体资源调配。面对“模型素材采集中存在困难”的风险时，果断组织评估和调整，确保方向正确。", 
                start_x+0.1, box_y+0.5, box_w-0.2, 1.4, 10, C_GRAY_TEXT)

    # 2. 技术破题
    start_x += box_w + gap
    box2 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(start_x), Inches(box_y), Inches(box_w), Inches(box_h))
    box2.fill.solid(); box2.fill.fore_color.rgb = C_LIGHT_BG; box2.line.color.rgb = C_BLUE
    add_textbox(s3, "🔧 技术破题：潘秦", start_x+0.1, box_y+0.1, box_w-0.2, 0.3, 11, C_BLUE, True)
    add_textbox(s3, "党员骨干主导模型实现，克服技术壁垒，誓保模型有效性达到99%。", 
                start_x+0.1, box_y+0.5, box_w-0.2, 1.4, 10, C_GRAY_TEXT)

    # 3. 党群联动
    start_x += box_w + gap
    box3 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(start_x), Inches(box_y), Inches(box_w), Inches(box_h))
    box3.fill.solid(); box3.fill.fore_color.rgb = C_LIGHT_BG; box3.line.color.rgb = C_GOLD
    add_textbox(s3, "🤝 党群联动机制 (1+N)", start_x+0.1, box_y+0.1, box_w-0.2, 0.3, 11, C_GOLD, True)
    add_textbox(s3, "带头人：党员骨干\n联动：张旭阳(预备党员/素材采集)、李鹏飞、丁灵(群众/系统开发)\n成效：形成党群协同联合攻关合力。", 
                start_x+0.1, box_y+0.5, box_w-0.2, 1.4, 10, C_GRAY_TEXT)

    s3_script = ('"关键在于构建了严密组织机制并发挥先锋模范作用。吴佳宇作为负责人，关键时刻评估调整方向。'
                 '潘秦主导模型实现，攻克技术壁垒。坚持党群共建，预备党员张旭阳冲锋在素材采集一线，群众骨干积极参与开发。'
                 '依靠党员领题、全员聚力，攻克了技术和管理难关。"').replace('\n', ' ')
    add_script(s3, s3_script)

    # ==================== 第 4 页：项目效益和成果 ====================
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s4, "项目效益和成果：质量跃升", "长效机制固化，创造实效", 4)

    # 左侧：数据对比
    add_textbox(s4, "📈 核心指标飞跃", 0.5, 1.4, 3, 0.3, 12, C_BLUE, True)
    
    # Before (修复处：直接赋值给变量，不使用_sp)
    oval_before = s4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(1.8), Inches(1.5), Inches(1.5))
    oval_before.fill.solid()
    oval_before.fill.fore_color.rgb = RGBColor(220, 220, 220)
    
    add_textbox(s4, "20%", 0.5, 2.2, 1.5, 0.5, 24, C_GRAY_TEXT, True, PP_ALIGN.CENTER)
    add_textbox(s4, "传统人工点检", 0.5, 2.7, 1.5, 0.3, 9, C_GRAY_TEXT, False, PP_ALIGN.CENTER)

    # Arrow
    arrow = s4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.2), Inches(2.3), Inches(0.8), Inches(0.4))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = C_GOLD

    # After (修复处：直接赋值给变量，不使用_sp)
    oval_after = s4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.2), Inches(1.8), Inches(1.5), Inches(1.5))
    oval_after.fill.solid()
    oval_after.fill.fore_color.rgb = C_BLUE
    
    add_textbox(s4, "90%+", 3.2, 2.2, 1.5, 0.5, 24, C_WHITE, True, PP_ALIGN.CENTER)
    add_textbox(s4, "AI侦测辅助监控", 3.2, 2.7, 1.5, 0.3, 9, C_WHITE, False, PP_ALIGN.CENTER)

    add_textbox(s4, "彻底解决了异常发现滞后的顽疾", 0.5, 3.5, 4.5, 0.3, 11, C_RED, True, PP_ALIGN.CENTER)

    # 右侧：成果列表
    list_x = 5.2
    list_y = 1.4
    
    items = [
        ("🔧 技术成果转化", "成功完成系统全面对接，预研落地了2个AI侦测辅助监控项目。"),
        ("🛡️ 管理价值深化", "有效弱化了此类‘管理X’对产品质量的影响度，实现质量跃升。"),
        ("🔄 机制固化", "改善成果已形成规范化操作流程，构建长效机制，杜绝问题反弹。")
    ]

    for title, desc in items:
        box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(list_x), Inches(list_y), Inches(4.4), Inches(0.9))
        box.fill.solid(); box.fill.fore_color.rgb = C_LIGHT_BG; box.line.color.rgb = C_BLUE
        
        add_textbox(s4, title, list_x+0.1, list_y+0.1, 4.2, 0.3, 11, C_BLUE, True)
        add_textbox(s4, desc, list_x+0.1, list_y+0.4, 4.2, 0.45, 10, C_GRAY_TEXT)
        list_y += 1.05

    s4_script = ('"最大价值在于实现了党建活力向发展动力的转化。核心指标实现质变，从不足20%提升至90%以上，实现实时精准监控。\n'
                 '不仅完成系统对接，落地2个AI项目，更将成果内化为规范化操作，构建了依靠数字化工具保持长期管理效果的制度体系。"').replace('\n', ' ')
    add_script(s4, s4_script)

    # ==================== 第 5 页：其他需说明事项 ====================
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s5, "其他需说明事项：问题剖析与长效机制", "持续改进，成果固化", 5)

    # 1. 风险与反思 (红色框)
    risk_box = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.375), Inches(1.4), Inches(9.25), Inches(1.0))
    risk_box.fill.solid(); risk_box.fill.fore_color.rgb = RGBColor(255, 235, 238)
    risk_box.line.color.rgb = C_RED; risk_box.line.dash_style = 1
    
    add_textbox(s5, "🛑 反思与差距", 0.5, 1.5, 3, 0.3, 12, C_RED, True)
    add_textbox(s5, "● 需警惕“两张皮”思想残余，重点关注模型持续优化。\n● 风险点：应对管理X随工艺流程调整而变化的风险。", 
                0.5, 1.8, 8.5, 0.5, 10, C_GRAY_TEXT)

    # 2. 后续重点工作
    add_textbox(s5, "📅 后续重点工作", 0.375, 2.6, 4, 0.3, 12, C_BLUE, True)
    work_box = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.375), Inches(2.95), Inches(4.5), Inches(1.0))
    work_box.fill.solid(); work_box.fill.fore_color.rgb = C_LIGHT_BG
    add_textbox(s5, "● 聚焦成果固化和赋能培训。\n● 截止时间：2025年11月25日前完成项目成果的现场落地应用。", 
                0.45, 3.0, 4.3, 0.9, 10, C_GRAY_TEXT)

    # 3. 长效机制
    add_textbox(s5, "🔄 长效机制 (PDCA)", 5.125, 2.6, 4, 0.3, 12, C_BLUE, True)
    mech_box = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.125), Inches(2.95), Inches(4.5), Inches(1.0))
    mech_box.fill.solid(); mech_box.fill.fore_color.rgb = C_LIGHT_BG
    add_textbox(s5, "● 持续推动AI辅助监控技术在公司其他业务场景的复制与推广。\n● 形成持续改进的闭环。", 
                5.2, 3.0, 4.3, 0.9, 10, C_GRAY_TEXT)

    s5_script = ('"结项之际坚持刀刃向内，深刻剖析差距。必须建立快速响应机制，优化模型以应对工艺变化。\n'
                 '下一步重点是成果巩固，确保11月25日前完成培训和移交。\n'
                 '我们将以本次攻坚为起点，深化‘党建+技术创新’，探索推广成功经验，提供坚实的红色动能。"').replace('\n', ' ')
    add_script(s5, s5_script)

    # 保存文件
    prs.save(FILENAME)
    print(f"成功生成演示文稿：{FILENAME}")

if __name__ == "__main__":
    create_presentation()