import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- 全局配置 ---
FILENAME = "2025年党员攻坚项目汇报_AI侦测项目3.pptx"

# 配色方案 (SCC Deep South Circuit Corporate Identity)
C_BLUE = RGBColor(0x00, 0x33, 0x66)      # 主色：深蓝
C_GOLD = RGBColor(0xB8, 0x86, 0x0B)      # 强调色：暗金
C_GRAY_TEXT = RGBColor(0x33, 0x33, 0x33) # 正文灰
C_LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)  # 浅灰背景 (用于卡片)
C_RED = RGBColor(0xD3, 0x2F, 0x2F)       # 警告/重点
C_GREEN = RGBColor(0x28, 0xA7, 0x45)     # 成功/完成
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BORDER = RGBColor(0xE0, 0xE0, 0xE0)    # 边框色

def create_presentation():
    # 1. 初始化 16:9 画布
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- 核心辅助函数：创建文本框 ---
    def add_textbox(slide, text, x, y, w, h, font_size=12, color=C_GRAY_TEXT, bold=False, align=None, bg_color=None, border_color=None):
        """快速创建一个文本框，支持背景色、边框和对齐"""
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        
        # 背景填充
        if bg_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg_color
        
        # 边框
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        
        if align:
            p.alignment = align
            
        run = p.add_run()
        run.text = text
        run.font.name = 'Microsoft YaHei'  # 强制雅黑
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        return shape

    # --- 核心辅助函数：绘制卡片 ---
    def add_card(slide, title, content, x, y, w, h, icon="📌"):
        """绘制带标题的卡片"""
        # 卡片背景
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_LIGHT_BG
        shape.line.color.rgb = C_BLUE
        shape.line.width = Pt(1.5)
        
        # 标题
        add_textbox(slide, f"{icon} {title}", x+0.1, y+0.1, w-0.2, 0.4, 13, C_BLUE, True, PP_ALIGN.LEFT)
        # 内容
        add_textbox(slide, content, x+0.1, y+0.6, w-0.2, h-0.7, 11, C_GRAY_TEXT, False, PP_ALIGN.LEFT)

    # --- 核心辅助函数：应用母版样式 ---
    def apply_master(slide, title_text, sub_text=None, page_num=None):
        # 1. 左上角 Logo
        logo = slide.shapes.add_textbox(Inches(0.4), Inches(0.4), Inches(3), Inches(0.5))
        logo.text_frame.text = "SCC深南电路"
        logo.text_frame.paragraphs[0].font.bold = True
        logo.text_frame.paragraphs[0].font.color.rgb = C_GOLD
        logo.text_frame.paragraphs[0].font.size = Pt(14)

        # 2. 标题与分割线
        add_textbox(slide, title_text, 0.4, 0.8, 8, 0.8, 24, C_BLUE, True)
        if sub_text:
             add_textbox(slide, f"|  {sub_text}", 4.5, 0.95, 8, 0.5, 16, C_GRAY_TEXT)
        
        # 分割线
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.5), Inches(12.53), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = C_GOLD
        line.line.fill.background()

        # 3. 页脚
        footer = slide.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(6), Inches(0.4))
        footer.text_frame.text = "2025年党员攻坚项目 | 南通深南党支部"
        footer.text_frame.paragraphs[0].font.size = Pt(9)
        footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

        if page_num:
            pg = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.5), Inches(0.4))
            pg.text_frame.text = str(page_num)
            pg.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            pg.text_frame.paragraphs[0].font.size = Pt(9)
            pg.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

    # --- 核心辅助函数：添加演讲稿框 ---
    def add_script(slide, script_text):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(5.8), Inches(12.53), Inches(1.1))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        bg.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        bg.line.dash_style = 1 
        
        add_textbox(slide, "演讲关键句：", 0.5, 5.9, 2, 0.3, 11, C_GOLD, True)
        content = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6))
        content.text_frame.word_wrap = True
        p = content.text_frame.paragraphs[0]
        p.text = script_text
        p.font.italic = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # ==================== 封面页 (Cover) ====================
    slide_cover = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 封面背景
    shape_bg = slide_cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.5), Inches(13.333), Inches(2))
    shape_bg.fill.solid(); shape_bg.fill.fore_color.rgb = C_BLUE; shape_bg.line.fill.background()
    
    logo_cover = slide_cover.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.5))
    logo_cover.text_frame.text = "SCC深南电路"
    logo_cover.text_frame.paragraphs[0].font.bold = True
    logo_cover.text_frame.paragraphs[0].font.color.rgb = C_GOLD
    logo_cover.text_frame.paragraphs[0].font.size = Pt(18)

    add_textbox(slide_cover, "AI侦测辅助监控项目", 1.0, 2.5, 10, 1.5, 48, C_BLUE, True)
    add_textbox(slide_cover, "2025年度党员攻坚项目汇报", 1.0, 3.5, 8, 0.8, 24, C_GRAY_TEXT)
    
    line_cover = slide_cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.3), Inches(1.5), Inches(0.1))
    line_cover.fill.solid(); line_cover.fill.fore_color.rgb = C_GOLD; line_cover.line.fill.background()

    add_textbox(slide_cover, "汇报支部：南通深南党支部", 1.0, 6.0, 6, 0.5, 16, C_WHITE, True)
    add_textbox(slide_cover, "项目负责人：吴佳宇", 1.0, 6.5, 6, 0.5, 16, C_WHITE, False)

    circle = slide_cover.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(1.5), Inches(5), Inches(5))
    circle.fill.background(); circle.line.color.rgb = RGBColor(0xF0, 0xF0, 0xF0); circle.line.width = Pt(20)

    # ==================== P1: 政治引领与顶层设计 ====================
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s1, "01 政治引领与顶层设计", "筑牢战斗堡垒，确立数字化总基调", 1)

    # 三张纵向卡片
    # 卡片1：理论
    add_card(s1, "理论武装", 
             "● 深入开展主题教育，强化政治思想引领。\n● 支部落实“三会一课”，专题研讨数字化转型战略。\n● 统一思想：以技术创新作为检验党建工作的试金石。", 
             0.5, 1.8, 3.8, 3.8)
    
    # 卡片2：方针
    add_card(s1, "顶层方针", 
             "● 确立“党建+业务”双轮驱动模式。\n● 战略承接：紧密对齐部门BSC战略指标。\n● 核心路径：数据驱动，智能辅助，降本增效。", 
             4.7, 1.8, 3.8, 3.8, "⚖️")

    # 卡片3：挑战 (红色边框强调)
    shape_c = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(1.8), Inches(3.8), Inches(3.8))
    shape_c.fill.solid(); shape_c.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
    shape_c.line.color.rgb = C_RED; shape_c.line.width = Pt(1.5)
    add_textbox(s1, "⚠️ 核心挑战与破局", 9.0, 1.9, 3.6, 0.4, 13, C_RED, True)
    add_textbox(s1, "● 现状痛点：传统人工点检有效性仅20%，异常发现滞后。\n● 攻坚目标：准确率>90%，实现AI智能化替代。\n● 紧迫性：需在8个月内完成从0到1的突破。", 
                9.0, 2.4, 3.6, 3.0, 11, C_GRAY_TEXT)

    add_script(s1, "“旗帜鲜明讲政治，凝心聚力促发展。面对人工点检效率低下的瓶颈，支部决定将AI侦测作为年度头号攻坚工程，发挥党组织把方向、管大局的作用。”")

    # ==================== P2: 组织攻坚与党群联动 ====================
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s2, "02 组织攻坚与党群联动", "党员领题突击，群众聚力攻坚", 2)

    # 顶部架构图：1+N 模式
    add_textbox(s2, "党建引领“1+N”攻坚战法", 0.5, 1.7, 5, 0.4, 14, C_BLUE, True)
    
    # 党员突击队 (左)
    party_box = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.2), Inches(3.5), Inches(3.2))
    party_box.fill.solid(); party_box.fill.fore_color.rgb = C_BLUE
    add_textbox(s2, "🔴 党员突击队", 0.6, 2.3, 3, 0.4, 14, C_WHITE, True)
    add_textbox(s2, "● 队长：吴佳宇\n   (统筹资源，风险把控)\n● 骨干：潘秦\n   (攻克模型算法壁垒)", 
                0.6, 2.9, 3.3, 2.0, 12, C_WHITE)

    # 链接箭头
    arr = s2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(3.5), Inches(1.0), Inches(0.6))
    arr.fill.solid(); arr.fill.fore_color.rgb = C_GOLD; arr.line.fill.background()
    add_textbox(s2, "带动/赋能", 4.2, 3.0, 1, 0.4, 10, C_GOLD, True, PP_ALIGN.CENTER)

    # 群众攻坚组 (中)
    mass_box = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.4), Inches(2.2), Inches(3.5), Inches(3.2))
    mass_box.fill.solid(); mass_box.fill.fore_color.rgb = C_LIGHT_BG; mass_box.line.color.rgb = C_BLUE
    add_textbox(s2, "🔵 群众攻坚组", 5.5, 2.3, 3, 0.4, 14, C_BLUE, True)
    add_textbox(s2, "● 先锋：张旭阳 (预备党员)\n   (发挥桥梁作用)\n● 成员：一线技术骨干\n   (负责现场采集/测试)", 
                5.5, 2.9, 3.3, 2.0, 12, C_GRAY_TEXT)

    # 联动机制 (右)
    mech_box = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(2.2), Inches(3.5), Inches(3.2))
    mech_box.fill.solid(); mech_box.fill.fore_color.rgb = C_LIGHT_BG; mech_box.line.color.rgb = C_GOLD
    add_textbox(s2, "⚙️ 联动机制", 9.3, 2.3, 3, 0.4, 14, C_GOLD, True)
    add_textbox(s2, "1. 党员立项，群众献策\n2. 党员破题，群众落地\n3. 师徒帮带，技能共享\n\n成果：解决现场工艺难点5项", 
                9.3, 2.9, 3.3, 2.0, 12, C_GRAY_TEXT)

    add_script(s2, "“一名党员就是一面旗帜。吴佳宇同志挂帅，潘秦同志攻坚，预备党员张旭阳同志带动群众广泛参与，形成了‘党员带头冲、群众跟着干’的生动局面。”")

    # ==================== P3: 重点项目全景推进表 ====================
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s3, "03 重点项目全景推进表", "挂图作战，节点管控，实绩说话", 3)

    # 表头
    headers = ["阶段/时间", "关键任务 (Milestone)", "牵头人", "当前状态", "难点与对策"]
    cols = [1.5, 3.5, 2.0, 2.0, 3.5]
    cx = 0.5
    for i, h in enumerate(headers):
        shape = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(1.8), Inches(cols[i]), Inches(0.5))
        shape.fill.solid(); shape.fill.fore_color.rgb = C_BLUE
        add_textbox(s3, h, cx, 1.9, cols[i], 0.4, 11, C_WHITE, True, PP_ALIGN.CENTER)
        cx += cols[i] + 0.1

    # 数据行
    rows = [
        ("3月15日", "数据夯实\n现状有效性分析", "吴佳宇", "✅ 已完成", "难点：历史数据缺失\n对策：全员补录"),
        ("4月10日", "硬件优化\n解决采集卡顿", "潘秦", "🚀 提前完成", "难点：设备兼容性\n对策：跨部门借调调试"),
        ("6月15日", "软件上线\n模型初步部署", "潘秦", "✅ 按期完成", "难点：误报率高\n对策：迭代训练算法"),
        ("9月20日", "系统互通\n现场实测联调", "张旭阳", "✅ 按期完成", "难点：接口协议不通\n对策：成立专项攻关组")
    ]

    cy = 2.4
    for r in rows:
        cx = 0.5
        for i, val in enumerate(r):
            # 状态颜色判断
            font_c = C_GRAY_TEXT
            if "提前" in val: font_c = C_GREEN
            elif "已完成" in val or "按期" in val: font_c = C_BLUE
            
            # 背景条
            shape = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(cy), Inches(cols[i]), Inches(0.7))
            shape.fill.solid(); shape.fill.fore_color.rgb = C_LIGHT_BG; shape.line.color.rgb = C_BORDER
            
            add_textbox(s3, val, cx, cy+0.1, cols[i], 0.6, 10, font_c, False, PP_ALIGN.CENTER)
            cx += cols[i] + 0.1
        cy += 0.8

    add_script(s3, "“我们坚持挂图作战，倒排工期。特别是在4月10日的硬件优化节点，攻坚小组通过连续奋战，比计划提前完成，为后续软件开发赢得了宝贵时间。”")

    # ==================== P4: 赋能增效与数据实绩 ====================
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s4, "04 赋能增效与数据实绩", "将党建活力转化为发展动力", 4)

    # 左侧：核心指标大字
    add_textbox(s4, "核心指标跃升", 0.5, 1.8, 4, 0.5, 14, C_BLUE, True)
    
    # 20% -> 90%
    box_Before = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.5), Inches(2.5), Inches(2.5))
    box_Before.fill.solid(); box_Before.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE); box_Before.line.fill.background()
    add_textbox(s4, "20%", 0.5, 3.2, 2.5, 1, 36, RGBColor(0x99, 0x99, 0x99), True, PP_ALIGN.CENTER)
    add_textbox(s4, "Before (人工)", 0.5, 4.2, 2.5, 0.5, 12, RGBColor(0x99, 0x99, 0x99), False, PP_ALIGN.CENTER)

    arrow_res = s4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.2), Inches(3.5), Inches(1.0), Inches(0.5))
    arrow_res.fill.solid(); arrow_res.fill.fore_color.rgb = C_GOLD

    box_After = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.4), Inches(2.5), Inches(2.5), Inches(2.5))
    box_After.fill.solid(); box_After.fill.fore_color.rgb = C_BLUE
    add_textbox(s4, ">90%", 4.4, 3.2, 2.5, 1, 36, C_WHITE, True, PP_ALIGN.CENTER)
    add_textbox(s4, "After (AI智能)", 4.4, 4.2, 2.5, 0.5, 12, C_WHITE, False, PP_ALIGN.CENTER)

    # 右侧：多维价值列表
    add_textbox(s4, "多维价值产出", 7.5, 1.8, 4, 0.5, 14, C_GOLD, True)
    
    vals = [
        ("💰 降本增效", "落地2个AI项目，大幅减少人力重复劳动，预计年化节省工时XXX小时。"),
        ("🛡️ 质量风控", "弱化人为不稳定因素，实现7*24小时无间断精准监控，异常拦截率提升至100%。"),
        ("📚 人才培养", "通过‘双培养’机制，将1名技术骨干发展为入党积极分子，锻炼了一支数字化先锋队。")
    ]
    
    vy = 2.5
    for vt, vd in vals:
        add_textbox(s4, vt, 7.5, vy, 5, 0.4, 12, C_BLUE, True)
        add_textbox(s4, vd, 7.5, vy+0.4, 5, 0.5, 11, C_GRAY_TEXT)
        vy += 1.0

    add_script(s4, "“不仅实现了点检准确率从20%到90%的质的飞跃，更重要的是，我们在实战中锻炼了队伍，证明了党建工作做实了就是生产力，做强了就是竞争力。”")

    # ==================== P5: 问题剖析与长效机制 ====================
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master(s5, "05 问题剖析与长效机制", "刀刃向内找差距，着眼长远建机制", 5)

    # 上半部分：反思 (红色虚线框)
    bg_reflect = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.5))
    bg_reflect.fill.background()
    bg_reflect.line.color.rgb = C_RED
    bg_reflect.line.dash_style = 3 # Dash
    
    add_textbox(s5, "🛑 存在不足与反思", 0.7, 1.9, 4, 0.4, 13, C_RED, True)
    add_textbox(s5, "1. 融合深度不够：部分环节仍存在‘两张皮’现象，技术语言与党建语言转换不够顺畅。\n2. 推广力度不足：目前仅在单工序应用，尚未形成全流程覆盖。", 
                0.7, 2.4, 11, 0.8, 11, C_GRAY_TEXT)

    # 下半部分：PDCA 循环计划 (三个圆圈)
    add_textbox(s5, "🔄 下一阶段部署 (PDCA)", 0.5, 3.6, 5, 0.4, 14, C_BLUE, True)

    steps = [
        ("持续迭代", "针对现场工艺变化\n优化AI模型参数"),
        ("全面移交", "11月25日前完成\n管理权平稳移交"),
        ("复制推广", "总结党建+技术经验\n推广至其他制程")
    ]
    
    sx = 1.0
    for i, (st, sd) in enumerate(steps):
        # 圆
        oval = s5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(sx), Inches(4.2), Inches(2.2), Inches(2.2))
        oval.fill.solid(); oval.fill.fore_color.rgb = C_WHITE
        oval.line.color.rgb = C_BLUE; oval.line.width = Pt(2.5)
        
        add_textbox(s5, st, sx, 4.8, 2.2, 0.4, 14, C_BLUE, True, PP_ALIGN.CENTER)
        add_textbox(s5, sd, sx, 5.3, 2.2, 0.8, 10, C_GRAY_TEXT, False, PP_ALIGN.CENTER)
        
        # 箭头
        if i < 2:
            arr_n = s5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(sx+2.4), Inches(5.1), Inches(0.8), Inches(0.4))
            arr_n.fill.solid(); arr_n.fill.fore_color.rgb = C_GOLD
            
        sx += 3.8

    add_script(s5, "“成绩属于过去。我们将坚持问题导向，刀刃向内，持续深化‘党建+技术’融合，致力将本次攻坚经验固化为长效机制，为公司高质量发展贡献持续的红色力量！”")

    # 保存文件
    try:
        prs.save(FILENAME)
        print(f"✅ 成功生成文件：{os.path.abspath(FILENAME)}")
        print("💡 提示：该版本已包含完整的党建逻辑闭环与视觉图表。")
    except Exception as e:
        print(f"❌ 生成失败：{e}")

if __name__ == "__main__":
    create_presentation()