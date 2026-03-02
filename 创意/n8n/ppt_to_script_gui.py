import streamlit as st
import json
import os
from pptx import Presentation  # 用于提取PPT文本
from openai import OpenAI  # 用于DeepSeek AI调用
import requests  # 用于模拟汇报（如果有后端API）
import time  # 用于模拟延迟，实际可移除

# 内部DeepSeek代理配置（根据你提供的）
client = OpenAI(
    api_key="0",  # 直接写 "0"
    base_url="http://ds.scc.com.cn/v1"  # 注意：去掉 /chat/completions，只留到 /v1
)
# 模拟后端汇报API（从n8n中提取，如果不需要，可移除）
#REPORT_URL = "http://10.30.43.199:8080/api/update_status"  # 从n8n JSON中

def extract_ppt_text(file_path):
    """提取PPT每页文本，类似n8n的PowerPoint to Text"""
    prs = Presentation(file_path)
    slides = {}
    for i, slide in enumerate(prs.slides, start=1):
        content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        content.append(run.text)
        slides[i] = "\n".join(content)
    return slides

def refine_core_theme(slide_number, content):
    """提炼核心主题和关键点，类似n8n的提炼核心-格式化输出"""
    prompt = f"""
你是一位资深企业内训讲师，擅长将 PPT 内容快速提炼成核心主题和讲解要点。

请阅读第 {slide_number} 页 PPT 内容:
{content}

请按以下格式严格输出 JSON：
{{
  "slideNumber": {slide_number},
  "coreTheme": "本页核心主题（一句话概括）",
  "keyPoints": ["关键点1", "关键点2", "关键点3"]
}}

要求：
1. 核心主题必须简洁明了，一句话概括本页内容。
2. 关键点列出 2-4 条最重要信息，可适度提炼和总结。
3. 输出必须严格为 JSON，不要加 Markdown 或多余文本。
"""
    response = client.chat.completions.create(
        model="ds-v3",  # 或 "ds-v3"
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7
    )
    try:
        return json.loads(response.choices[0].message.content)
    except:
        return {"slideNumber": slide_number, "coreTheme": "提取失败", "keyPoints": []}

def aggregate_syllabus(refined_slides):
    """聚合大纲，类似n8n的1. Code_聚合大纲"""
    syllabus_text = "以下是本课程的完整 PPT 大纲：\n"
    for slide in refined_slides:
        syllabus_text += f"[第 {slide['slideNumber']} 页] 主题：{slide['coreTheme']}\n      要点：{', '.join(slide['keyPoints'])}\n"
    return syllabus_text, len(refined_slides)

def director_plan(syllabus):
    """导演策划，类似n8n的2. Agent_导演策划"""
    prompt = f"""
你是一位经验极其丰富的演讲总策划（导演）。
你拥有全局视角，负责设计整个演讲的起承转合。

{syllabus}

任务：
请分析上述大纲的逻辑流，为每一页 PPT 生成一个【转场与逻辑指南】。

要求：
1. **关注衔接**：告诉撰稿人，这一页的开头如何承接上一页的内容，结尾如何引出下一页的内容。
2. **定义语气**：指出这一页应该用什么情绪（如：严肃强调、轻松举例、发人深省）。

请严格按照以下 JSON 格式输出一个数组（不要Markdown，只要纯 JSON）：
[
  {{
    "slideNumber": 1,
    "transition_instruction": "这是开场，要用...吸引注意力，引出主题...",
    "mood": "激昂"
  }},
  ...
]
"""
    response = client.chat.completions.create(
        model="ds-v3",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7
    )
    try:
        return json.loads(response.choices[0].message.content)
    except:
        return []

def merge_guides(original_slides, refined_slides, director_guides):
    """合并指南，类似n8n的3. Code_合并指南"""
    merged = []
    for refined in refined_slides:
        slide_num = refined['slideNumber']
        guide = next((g for g in director_guides if g['slideNumber'] == slide_num), {})
        original_content = original_slides.get(slide_num, "未找到")
        merged.append({
            "slideNumber": slide_num,
            "original_content": original_content,
            "coreTheme": refined.get("coreTheme", ""),
            "keyPoints": refined.get("keyPoints", []),
            "director_instruction": guide.get("transition_instruction", "本页无特殊转场要求，请自然衔接。"),
            "mood": guide.get("mood", "自然")
        })
    return merged

def generate_draft(slide_data):
    """生成初稿，类似n8n的1. Writer_初稿生成"""
    prompt = f"""
你是一位专业的课程内容撰稿人。

【当前任务】：撰写第 {slide_data['slideNumber']} 页的数字人讲稿。

【输入素材】：
1. 核心主题：{slide_data['coreTheme']}
2. 关键要点：
{'\n'.join(slide_data['keyPoints'])}
3. **PPT 原文参考（用于补充细节与数据）**：
{slide_data['original_content']}

【🔴 导演的转场指令】：
{slide_data['director_instruction']}
(语气要求：{slide_data['mood']})

【撰写要求】：
1. **衔接**：必须严格执行“导演的转场指令”，确保上下文连贯。
2. **内容**：以“核心主题”为主骨架，从“PPT原文参考”中提取具体案例、数据或描述性语句来丰富内容，**严禁瞎编数据**。
3. **口语化**：生成适合口播的初稿。
"""
    response = client.chat.completions.create(
        model="ds-v3",
        messages=[{"role": "system", "content": "你是一个只负责生成初稿的AI助手。"}, {"role": "user", "content": prompt}]
    )
    return response.choices[0].message.content

def criticize_draft(draft, slide_number):
    """毒舌教官，类似n8n的2. Critic_毒舌教官"""
    prompt = f"""
你是一位极其严格的“数字人演讲教练”。你的任务不是改写内容，而是**审视**上一步生成的讲稿，并提出尖锐的修改意见。

【待评审初稿】：
{draft}

请从以下维度进行“反思”和批判：
1. **口语化程度**：是否有太多书面语（如“综上所述”/“由此可见”）？数字人应该像真人聊天一样说话。
2. **呼吸感**：句子是否太长？数字人需要换气，长难句会导致语音合成不自然。
3. **对象感**：是否像在对着空气念经？需要有“你”和“我”的交流感，要像对着一位坐在对面的新员工说话。

输出格式要求：
请列出 3-5 条具体的修改建议（Critique），必须用纯文本列表形式，不要输出 JSON。
"""
    response = client.chat.completions.create(
        model="ds-v3",
        messages=[{"role": "system", "content": "你是一个毒舌但专业的演讲教练。只输出批评意见，不输出修改后的文案。"}, {"role": "user", "content": prompt}]
    )
    return response.choices[0].message.content

def refine_final(draft, critiques, slide_number):
    """最终精修，类似n8n的3. Editor_最终精修"""
    prompt = f"""
你是一位顶级的演讲稿精修师。

输入信息：
1. [初稿内容]：
{draft}

2. [修改建议]：
{critiques}

任务：
请根据 [修改建议] 对 [初稿] 进行深度重写。

最终输出要求：
1. **极度口语化**：使用短句、反问句、感叹句。
2. **标注停顿**：在需要数字人明显停顿的地方，可以使用逗号或空格来控制节奏。
3. **最终格式**：只输出最终的演讲稿文本，不要包含“根据建议修改如下”之类的废话。

现在，请输出第 {slide_number} 页的最终讲稿：
"""
    response = client.chat.completions.create(
        model="ds-v3",
        messages=[{"role": "system", "content": "你是一个金牌编辑，负责输出最终完美的口语化脚本。"}, {"role": "user", "content": prompt}]
    )
    return response.choices[0].message.content

def report_progress(task_id, status, data=None):
    """模拟汇报进度，类似n8n的HTTP Request节点"""
    payload = {"task_id": task_id, "status": status}
    if data:
        payload["data"] = json.dumps(data)
    try:
        requests.post(REPORT_URL, json=payload)
    except:
        pass  # 如果无后端，忽略

# Streamlit GUI
st.title("PPT转数字人脚本工具")

# 上传文件（支持拖拽）
uploaded_file = st.file_uploader("拖拽或选择PPT文件上传", type=["pptx"])

if uploaded_file:
    # 保存上传文件
    file_path = os.path.join("temp", uploaded_file.name)
    os.makedirs("temp", exist_ok=True)
    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    
    task_id = uploaded_file.name  # 简单用文件名作为task_id
    
    with st.status("处理中...", expanded=True) as status:
        st.write("步骤1: 提取PPT文本...")
        original_slides = extract_ppt_text(file_path)
        report_progress(task_id, "paginated", {"total_slides": len(original_slides)})
        time.sleep(1)  # 模拟延迟
        
        st.write("步骤2: 提炼每页核心主题...")
        refined_slides = []
        progress_bar = st.progress(0)
        for i, (slide_num, content) in enumerate(original_slides.items()):
            refined = refine_core_theme(slide_num, content)
            refined_slides.append(refined)
            progress_bar.progress((i + 1) / len(original_slides))
        
        st.write("步骤3: 聚合大纲...")
        syllabus, total_slides = aggregate_syllabus(refined_slides)
        report_progress(task_id, "syllabus_ready", {"syllabus": syllabus})
        
        st.write("步骤4: 导演策划...")
        director_guides = director_plan(syllabus)
        report_progress(task_id, "director_finished", {"msg": "导演策划已完成"})
        
        st.write("步骤5: 合并指南并生成脚本...")
        merged_slides = merge_guides(original_slides, refined_slides, director_guides)
        
        final_scripts = {}
        for slide in merged_slides:
            draft = generate_draft(slide)
            critiques = criticize_draft(draft, slide['slideNumber'])
            final = refine_final(draft, critiques, slide['slideNumber'])
            final_scripts[slide['slideNumber']] = final
            report_progress(task_id, "page_generated", {
                "slide_number": slide['slideNumber'],
                "content": final
            })
            st.write(f"第 {slide['slideNumber']} 页完成")
        
        status.update(label="处理完成！", state="complete")
        report_progress(task_id, "completed", {"final_result": final_scripts})
    
    # 显示结果并提供下载
    st.subheader("最终结果")
    st.json(final_scripts)
    
    # 下载按钮
    result_json = json.dumps(final_scripts, ensure_ascii=False, indent=4)
    st.download_button(
        label="下载结果 (JSON)",
        data=result_json,
        file_name="final_scripts.json",
        mime="application/json"
    )

    # 清理临时文件
    os.remove(file_path)