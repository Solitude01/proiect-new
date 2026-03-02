#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Office文档敏感关键词检测工具
支持检查 Word (.docx)、Excel (.xlsx)、PPT (.pptx) 文档中的敏感关键词
"""

import os
import sys
import json
import datetime
import threading
import queue
from pathlib import Path
from typing import List, Dict, Set, Tuple
from concurrent.futures import ThreadPoolExecutor, as_completed

import customtkinter as ctk
from tkinter import filedialog, messagebox
from tkinterdnd2 import DND_FILES, TkinterDnD

# Office 文档处理库
from docx import Document as WordDocument
from openpyxl import load_workbook
from pptx import Presentation


# ==================== 配置管理模块 ====================
class ConfigManager:
    """配置管理器，负责保存和加载用户设置"""

    def __init__(self, config_file="config/settings.json"):
        self.config_file = config_file
        self.default_config = {
            "default_keywords": {
                "企密A": True,
                "企密AA": True,
                "企密AAA": True
            },
            "custom_keywords": [],
            "check_content": False,
            "last_directory": ""
        }
        self._ensure_config_dir()

    def _ensure_config_dir(self):
        """确保配置目录存在"""
        os.makedirs(os.path.dirname(self.config_file), exist_ok=True)

    def load(self) -> dict:
        """加载配置"""
        try:
            if os.path.exists(self.config_file):
                with open(self.config_file, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                    # 合并默认配置（处理新增配置项）
                    for key in self.default_config:
                        if key not in config:
                            config[key] = self.default_config[key]
                    return config
        except Exception as e:
            print(f"加载配置失败: {e}")
        return self.default_config.copy()

    def save(self, config: dict):
        """保存配置"""
        try:
            self._ensure_config_dir()
            with open(self.config_file, 'w', encoding='utf-8') as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"保存配置失败: {e}")


# ==================== 文档检查核心逻辑 ====================
class DocumentChecker:
    """文档检查器，负责检查各类Office文档中的关键词"""

    def __init__(self, keywords: List[str], check_content: bool = False):
        self.keywords_original = keywords  # 保存原始关键词（用于显示）
        self.keywords_lower = [kw.lower() for kw in keywords]  # 转小写用于忽略大小写匹配
        self.check_content = check_content

    @staticmethod
    def _get_excel_column_name(col_num: int) -> str:
        """将列号转换为Excel列名（A, B, ..., Z, AA, AB, ...）"""
        result = ""
        while col_num > 0:
            col_num -= 1
            result = chr(65 + col_num % 26) + result
            col_num //= 26
        return result

    def check_word(self, file_path: str) -> List[Dict]:
        """检查Word文档"""
        results = []
        try:
            doc = WordDocument(file_path)

            # 检查页眉
            for section in doc.sections:
                header = section.header
                header_text = self._get_header_footer_text(header)
                matches = self._find_matches(header_text)
                for keyword in matches:
                    results.append({
                        "keyword": keyword,
                        "location": "页眉",
                        "context": header_text[:100],
                        "type": "header_footer"  # 标识为页眉页脚
                    })

            # 检查页脚
            for section in doc.sections:
                footer = section.footer
                footer_text = self._get_header_footer_text(footer)
                matches = self._find_matches(footer_text)
                for keyword in matches:
                    results.append({
                        "keyword": keyword,
                        "location": "页脚",
                        "context": footer_text[:100],
                        "type": "header_footer"  # 标识为页眉页脚
                    })

            # 检查正文（如果启用）
            if self.check_content:
                for para_idx, para in enumerate(doc.paragraphs):
                    matches = self._find_matches(para.text)
                    for keyword in matches:
                        results.append({
                            "keyword": keyword,
                            "location": f"正文-段落{para_idx + 1}",
                            "context": para.text[:100],
                            "type": "content"  # 标识为正文
                        })

                # 检查表格
                for table_idx, table in enumerate(doc.tables):
                    for row_idx, row in enumerate(table.rows):
                        for cell_idx, cell in enumerate(row.cells):
                            matches = self._find_matches(cell.text)
                            for keyword in matches:
                                results.append({
                                    "keyword": keyword,
                                    "location": f"正文-表格{table_idx + 1}-行{row_idx + 1}-列{cell_idx + 1}",
                                    "context": cell.text[:100],
                                    "type": "content"  # 标识为正文
                                })

        except Exception as e:
            raise Exception(f"Word文档处理错误: {str(e)}")

        return results

    def check_excel(self, file_path: str) -> List[Dict]:
        """检查Excel文档"""
        results = []
        try:
            wb = load_workbook(file_path, data_only=True)

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]

                # 检查页眉页脚
                if hasattr(sheet, 'HeaderFooter'):
                    hf = sheet.HeaderFooter
                    # 检查页眉
                    header_parts = [hf.left_header, hf.center_header, hf.right_header]
                    header_text = ' '.join([str(p) for p in header_parts if p])
                    matches = self._find_matches(header_text)
                    for keyword in matches:
                        results.append({
                            "keyword": keyword,
                            "location": f"工作表'{sheet_name}'-页眉",
                            "context": header_text[:100],
                            "type": "header_footer"
                        })

                    # 检查页脚
                    footer_parts = [hf.left_footer, hf.center_footer, hf.right_footer]
                    footer_text = ' '.join([str(p) for p in footer_parts if p])
                    matches = self._find_matches(footer_text)
                    for keyword in matches:
                        results.append({
                            "keyword": keyword,
                            "location": f"工作表'{sheet_name}'-页脚",
                            "context": footer_text[:100],
                            "type": "header_footer"
                        })

                # 检查单元格内容（如果启用）
                if self.check_content:
                    for row_idx, row in enumerate(sheet.iter_rows(values_only=True), 1):
                        for col_idx, cell_value in enumerate(row, 1):
                            if cell_value:
                                cell_text = str(cell_value)
                                matches = self._find_matches(cell_text)
                                for keyword in matches:
                                    col_name = self._get_excel_column_name(col_idx)
                                    results.append({
                                        "keyword": keyword,
                                        "location": f"工作表'{sheet_name}'-单元格{col_name}{row_idx}",
                                        "context": cell_text[:100],
                                        "type": "content"
                                    })

            wb.close()

        except Exception as e:
            raise Exception(f"Excel文档处理错误: {str(e)}")

        return results

    def check_ppt(self, file_path: str) -> List[Dict]:
        """检查PowerPoint文档"""
        results = []
        try:
            prs = Presentation(file_path)

            # PPT没有真正的页眉页脚，所有内容都算作正文
            # 只有在启用"检测正文"时才检查
            if self.check_content:
                # 检查每张幻灯片的形状和表格
                for slide_idx, slide in enumerate(prs.slides, 1):
                    # 检查幻灯片中的所有形状
                    for shape in slide.shapes:
                        # 检查文本框
                        if hasattr(shape, "text"):
                            text = shape.text
                            matches = self._find_matches(text)
                            for keyword in matches:
                                location = f"幻灯片{slide_idx}"
                                if shape.name:
                                    location += f"-{shape.name}"
                                results.append({
                                    "keyword": keyword,
                                    "location": location,
                                    "context": text[:100],
                                    "type": "content"  # PPT幻灯片内容算作正文
                                })

                        # 检查表格
                        if shape.has_table:
                            table = shape.table
                            for row_idx, row in enumerate(table.rows):
                                for cell_idx, cell in enumerate(row.cells):
                                    text = cell.text
                                    matches = self._find_matches(text)
                                    for keyword in matches:
                                        results.append({
                                            "keyword": keyword,
                                            "location": f"幻灯片{slide_idx}-表格-行{row_idx + 1}-列{cell_idx + 1}",
                                            "context": text[:100],
                                            "type": "content"  # 表格内容也算作正文
                                        })

                # 检查备注
                for slide_idx, slide in enumerate(prs.slides, 1):
                    if slide.has_notes_slide:
                        notes_text = slide.notes_slide.notes_text_frame.text
                        matches = self._find_matches(notes_text)
                        for keyword in matches:
                            results.append({
                                "keyword": keyword,
                                "location": f"幻灯片{slide_idx}-备注",
                                "context": notes_text[:100],
                                "type": "content"  # 备注算作正文
                            })

        except Exception as e:
            raise Exception(f"PPT文档处理错误: {str(e)}")

        return results

    def _get_header_footer_text(self, header_footer) -> str:
        """获取Word页眉页脚文本"""
        text_parts = []
        for para in header_footer.paragraphs:
            text_parts.append(para.text)
        for table in header_footer.tables:
            for row in table.rows:
                for cell in row.cells:
                    text_parts.append(cell.text)
        return ' '.join(text_parts)

    def _find_matches(self, text: str) -> Set[str]:
        """在文本中查找匹配的关键词（优先匹配最长的关键词，避免重复匹配）"""
        if not text:
            return set()

        text_lower = text.lower()
        matches = set()
        matched_positions = []  # 记录已匹配的位置范围

        # 按关键词长度从长到短排序，优先匹配长关键词
        sorted_keywords = sorted(
            enumerate(self.keywords_lower),
            key=lambda x: len(x[1]),
            reverse=True
        )

        for i, keyword_lower in sorted_keywords:
            # 查找所有出现位置
            start = 0
            while True:
                pos = text_lower.find(keyword_lower, start)
                if pos == -1:
                    break

                end = pos + len(keyword_lower)

                # 检查这个位置是否已经被更长的关键词匹配过
                is_overlapping = False
                for matched_start, matched_end in matched_positions:
                    if not (end <= matched_start or pos >= matched_end):
                        is_overlapping = True
                        break

                if not is_overlapping:
                    # 返回原始关键词（保持大小写）
                    matches.add(self.keywords_original[i])
                    matched_positions.append((pos, end))

                start = pos + 1

        return matches


# ==================== 日志管理器 ====================
class LogManager:
    """日志管理器，负责日志的显示和保存"""

    def __init__(self, log_dir="log"):
        self.log_dir = log_dir
        self.log_file = None
        self.log_buffer = []
        self._ensure_log_dir()

    def _ensure_log_dir(self):
        """确保日志目录存在"""
        os.makedirs(self.log_dir, exist_ok=True)

    def start_logging(self):
        """开始新的日志记录"""
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M")
        self.log_file = os.path.join(self.log_dir, f"checklog_{timestamp}.txt")
        self.log_buffer = []

    def log(self, message: str, level="INFO") -> str:
        """记录日志"""
        timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        log_entry = f"[{timestamp}] {message}"

        self.log_buffer.append(log_entry)

        # 写入文件
        if self.log_file:
            try:
                with open(self.log_file, 'a', encoding='utf-8') as f:
                    f.write(log_entry + '\n')
            except Exception as e:
                print(f"写入日志文件失败: {e}")

        return log_entry

    def save_to_file(self, file_path: str):
        """保存日志到指定文件"""
        try:
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(self.log_buffer))
            return True
        except Exception as e:
            print(f"保存日志失败: {e}")
            return False


# ==================== GUI主界面 ====================
class DocumentCheckerGUI(TkinterDnD.Tk):
    """主GUI界面"""

    def __init__(self):
        super().__init__()

        # 窗口设置
        self.title("Office文档敏感关键词检测工具")
        self.geometry("1000x700")

        # 设置主题
        ctk.set_appearance_mode("light")
        ctk.set_default_color_theme("blue")

        # 初始化管理器
        self.config_manager = ConfigManager()
        self.log_manager = LogManager()

        # 加载配置
        self.settings = self.config_manager.load()

        # 数据变量
        self.file_items = {}  # 待检查的文件：{路径: {"checked": BooleanVar, "name": 文件名}}
        self.default_keyword_vars = {}  # 默认关键词复选框变量
        self.custom_keywords = self.settings.get("custom_keywords", []).copy()
        self.check_content_var = ctk.BooleanVar(value=self.settings.get("check_content", False))
        self.is_checking = False

        # 检查结果
        self.check_results = []

        # 创建界面
        self._create_widgets()

        # 设置拖拽
        self._setup_drag_drop()

    def _create_widgets(self):
        """创建界面组件"""

        # ===== 左侧面板 - 关键词设置 =====
        left_frame = ctk.CTkFrame(self, width=280)
        left_frame.pack(side="left", fill="y", padx=10, pady=10)
        left_frame.pack_propagate(False)

        # 标题
        title_label = ctk.CTkLabel(
            left_frame,
            text="敏感词检测工具",
            font=ctk.CTkFont(size=18, weight="bold")
        )
        title_label.pack(pady=(10, 15))

        # 默认关键词区域
        default_kw_frame = ctk.CTkFrame(left_frame)
        default_kw_frame.pack(fill="x", padx=10, pady=10)

        ctk.CTkLabel(
            default_kw_frame,
            text="默认关键词",
            font=ctk.CTkFont(size=14, weight="bold")
        ).pack(pady=5)

        default_keywords = ["企密A", "企密AA", "企密AAA"]
        for kw in default_keywords:
            var = ctk.BooleanVar(value=self.settings["default_keywords"].get(kw, True))
            self.default_keyword_vars[kw] = var
            checkbox = ctk.CTkCheckBox(default_kw_frame, text=kw, variable=var)
            checkbox.pack(anchor="w", padx=20, pady=2)

        # 自定义关键词区域
        custom_kw_frame = ctk.CTkFrame(left_frame)
        custom_kw_frame.pack(fill="both", expand=True, padx=10, pady=10)

        ctk.CTkLabel(
            custom_kw_frame,
            text="自定义关键词",
            font=ctk.CTkFont(size=14, weight="bold")
        ).pack(pady=5)

        # 自定义关键词输入
        input_frame = ctk.CTkFrame(custom_kw_frame)
        input_frame.pack(fill="x", padx=5, pady=5)

        self.custom_kw_entry = ctk.CTkEntry(input_frame, placeholder_text="输入关键词")
        self.custom_kw_entry.pack(side="left", fill="x", expand=True, padx=2)
        self.custom_kw_entry.bind("<Return>", lambda e: self._add_custom_keyword())  # 支持Enter键

        add_kw_btn = ctk.CTkButton(
            input_frame,
            text="添加",
            width=60,
            command=self._add_custom_keyword
        )
        add_kw_btn.pack(side="right", padx=2)

        # 自定义关键词列表
        self.custom_kw_listbox = ctk.CTkTextbox(custom_kw_frame, height=100)
        self.custom_kw_listbox.pack(fill="both", expand=True, padx=5, pady=5)
        self._update_custom_keyword_list()

        delete_kw_btn = ctk.CTkButton(
            custom_kw_frame,
            text="删除选中",
            command=self._delete_custom_keyword,
            fg_color="red"
        )
        delete_kw_btn.pack(fill="x", padx=5, pady=2)

        # 检测选项
        options_frame = ctk.CTkFrame(left_frame)
        options_frame.pack(fill="x", padx=10, pady=10)

        ctk.CTkLabel(
            options_frame,
            text="检测选项",
            font=ctk.CTkFont(size=14, weight="bold")
        ).pack(pady=5)

        self.check_content_checkbox = ctk.CTkCheckBox(
            options_frame,
            text="检测正文关键词",
            variable=self.check_content_var
        )
        self.check_content_checkbox.pack(anchor="w", padx=20, pady=2)

        # 操作按钮
        self.start_btn = ctk.CTkButton(
            left_frame,
            text="开始检查",
            command=self._start_check,
            height=40,
            font=ctk.CTkFont(size=16, weight="bold"),
            fg_color="green"
        )
        self.start_btn.pack(fill="x", padx=10, pady=10)

        # ===== 中间面板 - 文件管理 =====
        middle_frame = ctk.CTkFrame(self, width=380)
        middle_frame.pack(side="left", fill="both", expand=False, padx=(0, 10), pady=10)
        middle_frame.pack_propagate(False)

        # 文件管理标题
        ctk.CTkLabel(
            middle_frame,
            text="文件管理",
            font=ctk.CTkFont(size=16, weight="bold")
        ).pack(pady=(10, 10))

        # 文件选择按钮区
        file_btn_frame = ctk.CTkFrame(middle_frame)
        file_btn_frame.pack(fill="x", padx=10, pady=5)

        self.select_files_btn = ctk.CTkButton(
            file_btn_frame,
            text="选择文件",
            command=self._select_files,
            width=120
        )
        self.select_files_btn.pack(side="left", padx=2)

        self.select_folder_btn = ctk.CTkButton(
            file_btn_frame,
            text="选择文件夹",
            command=self._select_folder,
            width=120
        )
        self.select_folder_btn.pack(side="left", padx=2)

        self.clear_files_btn = ctk.CTkButton(
            file_btn_frame,
            text="清空列表",
            command=self._clear_files,
            fg_color="gray",
            width=100
        )
        self.clear_files_btn.pack(side="left", padx=2)

        # 批量添加路径区域
        batch_frame = ctk.CTkFrame(middle_frame)
        batch_frame.pack(fill="x", padx=10, pady=5)

        ctk.CTkLabel(
            batch_frame,
            text="批量添加路径（每行一个）：",
            font=ctk.CTkFont(size=12, weight="bold")
        ).pack(anchor="w", padx=5, pady=(5, 2))

        # 路径输入文本框
        self.path_input_text = ctk.CTkTextbox(batch_frame, height=80, font=ctk.CTkFont(size=11))
        self.path_input_text.pack(fill="x", padx=5, pady=5)

        # 添加路径按钮
        add_paths_btn = ctk.CTkButton(
            batch_frame,
            text="添加路径",
            command=self._add_paths_from_text,
            fg_color="blue",
            width=100
        )
        add_paths_btn.pack(pady=5)

        # 文件统计标签
        self.file_count_label = ctk.CTkLabel(
            middle_frame,
            text="已添加: 0 个文件 | 已选中: 0 个",
            font=ctk.CTkFont(size=12)
        )
        self.file_count_label.pack(pady=5)

        # 文件列表区域
        file_list_frame = ctk.CTkFrame(middle_frame)
        file_list_frame.pack(fill="both", expand=True, padx=10, pady=5)

        ctk.CTkLabel(
            file_list_frame,
            text="文件列表（勾选要检查的文件）：",
            font=ctk.CTkFont(size=12, weight="bold")
        ).pack(anchor="w", padx=5, pady=5)

        # 全选/全不选按钮
        select_btn_frame = ctk.CTkFrame(file_list_frame, fg_color="transparent")
        select_btn_frame.pack(fill="x", padx=5, pady=2)

        self.select_all_btn = ctk.CTkButton(
            select_btn_frame,
            text="全选",
            command=self._select_all_files,
            width=80,
            height=25,
            font=ctk.CTkFont(size=11)
        )
        self.select_all_btn.pack(side="left", padx=2)

        self.deselect_all_btn = ctk.CTkButton(
            select_btn_frame,
            text="全不选",
            command=self._deselect_all_files,
            width=80,
            height=25,
            font=ctk.CTkFont(size=11),
            fg_color="gray"
        )
        self.deselect_all_btn.pack(side="left", padx=2)

        # 滚动文件列表
        self.file_list_scroll = ctk.CTkScrollableFrame(file_list_frame, label_text="")
        self.file_list_scroll.pack(fill="both", expand=True, padx=5, pady=5)

        # ===== 右侧面板 - 日志显示 =====
        right_frame = ctk.CTkFrame(self)
        right_frame.pack(side="right", fill="both", expand=True, padx=(0, 10), pady=10)

        # 日志区域标题
        log_title_frame = ctk.CTkFrame(right_frame)
        log_title_frame.pack(fill="x", padx=10, pady=(10, 5))

        ctk.CTkLabel(
            log_title_frame,
            text="检查日志",
            font=ctk.CTkFont(size=16, weight="bold")
        ).pack(side="left")

        # 日志操作按钮
        log_btn_frame = ctk.CTkFrame(log_title_frame, fg_color="transparent")
        log_btn_frame.pack(side="right")

        self.save_log_btn = ctk.CTkButton(
            log_btn_frame,
            text="保存日志",
            width=100,
            command=self._save_log
        )
        self.save_log_btn.pack(side="left", padx=5)

        self.export_report_btn = ctk.CTkButton(
            log_btn_frame,
            text="导出报告",
            width=100,
            command=self._export_report
        )
        self.export_report_btn.pack(side="left", padx=5)

        clear_log_btn = ctk.CTkButton(
            log_btn_frame,
            text="清空日志",
            width=100,
            command=self._clear_log,
            fg_color="gray"
        )
        clear_log_btn.pack(side="left", padx=5)

        # 日志文本框
        self.log_text = ctk.CTkTextbox(right_frame, font=ctk.CTkFont(size=12))
        self.log_text.pack(fill="both", expand=True, padx=10, pady=5)

        # 进度条和状态栏
        status_frame = ctk.CTkFrame(right_frame)
        status_frame.pack(fill="x", padx=10, pady=(5, 10))

        self.progress_bar = ctk.CTkProgressBar(status_frame)
        self.progress_bar.pack(fill="x", pady=5)
        self.progress_bar.set(0)

        self.status_label = ctk.CTkLabel(
            status_frame,
            text="就绪",
            font=ctk.CTkFont(size=12)
        )
        self.status_label.pack()

    def _setup_drag_drop(self):
        """设置拖拽功能"""
        self.drop_target_register(DND_FILES)
        self.dnd_bind('<<Drop>>', self._on_drop)

    def _on_drop(self, event):
        """处理拖拽文件"""
        files = self.tk.splitlist(event.data)
        self._add_files(files)

    def _select_files(self):
        """选择文件"""
        files = filedialog.askopenfilenames(
            title="选择文件",
            filetypes=[
                ("Office文档", "*.docx *.xlsx *.pptx"),
                ("Word文档", "*.docx"),
                ("Excel文档", "*.xlsx"),
                ("PPT文档", "*.pptx"),
                ("所有文件", "*.*")
            ],
            initialdir=self.settings.get("last_directory", "")
        )
        if files:
            self._add_files(files)
            # 保存最后使用的目录
            self.settings["last_directory"] = os.path.dirname(files[0])
            self.config_manager.save(self.settings)

    def _select_folder(self):
        """选择文件夹"""
        folder = filedialog.askdirectory(
            title="选择文件夹",
            initialdir=self.settings.get("last_directory", "")
        )
        if folder:
            self._add_files([folder])
            self.settings["last_directory"] = folder
            self.config_manager.save(self.settings)

    def _add_files(self, paths):
        """添加文件到列表"""
        valid_extensions = {'.docx', '.xlsx', '.pptx'}
        added_count = 0

        for path in paths:
            path = path.strip('{}').strip('"').strip("'")  # 移除可能的大括号和引号

            if os.path.isfile(path):
                # 单个文件
                if Path(path).suffix.lower() in valid_extensions:
                    if path not in self.file_items:
                        # 添加到文件列表，默认不勾选
                        var = ctk.BooleanVar(value=False)
                        self.file_items[path] = {
                            "checked": var,
                            "name": os.path.basename(path)
                        }
                        added_count += 1

            elif os.path.isdir(path):
                # 文件夹，递归查找
                for root, dirs, files in os.walk(path):
                    for file in files:
                        if Path(file).suffix.lower() in valid_extensions:
                            file_path = os.path.join(root, file)
                            if file_path not in self.file_items:
                                var = ctk.BooleanVar(value=False)
                                self.file_items[file_path] = {
                                    "checked": var,
                                    "name": os.path.basename(file_path)
                                }
                                added_count += 1

        if added_count > 0:
            self._update_file_list_display()
            self._update_file_count()
            self._log(f"已添加 {added_count} 个文件")

    def _add_paths_from_text(self):
        """从文本框批量添加路径"""
        text_content = self.path_input_text.get("1.0", "end").strip()
        if not text_content:
            messagebox.showwarning("提示", "请输入文件路径")
            return

        # 按行分割路径
        paths = []
        for line in text_content.split('\n'):
            line = line.strip()
            if line:  # 忽略空行
                paths.append(line)

        if paths:
            self._add_files(paths)
            # 清空输入框
            self.path_input_text.delete("1.0", "end")

    def _update_file_list_display(self):
        """更新文件列表显示"""
        # 清空现有显示
        for widget in self.file_list_scroll.winfo_children():
            widget.destroy()

        # 按文件名排序显示
        sorted_files = sorted(self.file_items.items(), key=lambda x: x[1]["name"])

        for file_path, file_info in sorted_files:
            # 创建一个frame包含复选框
            file_frame = ctk.CTkFrame(self.file_list_scroll, fg_color="transparent")
            file_frame.pack(fill="x", pady=2, padx=5)

            # 复选框
            checkbox = ctk.CTkCheckBox(
                file_frame,
                text=file_info["name"],
                variable=file_info["checked"],
                font=ctk.CTkFont(size=11),
                command=self._update_file_count
            )
            checkbox.pack(side="left", fill="x", expand=True)

            # 删除按钮
            delete_btn = ctk.CTkButton(
                file_frame,
                text="×",
                width=30,
                height=25,
                command=lambda fp=file_path: self._remove_single_file(fp),
                fg_color="red",
                hover_color="darkred",
                font=ctk.CTkFont(size=14, weight="bold")
            )
            delete_btn.pack(side="right")

    def _remove_single_file(self, file_path):
        """删除单个文件"""
        if file_path in self.file_items:
            del self.file_items[file_path]
            self._update_file_list_display()
            self._update_file_count()
            self._log(f"已移除文件: {os.path.basename(file_path)}")

    def _select_all_files(self):
        """全选所有文件"""
        for file_info in self.file_items.values():
            file_info["checked"].set(True)
        self._update_file_count()

    def _deselect_all_files(self):
        """取消选择所有文件"""
        for file_info in self.file_items.values():
            file_info["checked"].set(False)
        self._update_file_count()

    def _clear_files(self):
        """清空文件列表"""
        self.file_items.clear()
        self._update_file_list_display()
        self._update_file_count()
        self._log("已清空文件列表")

    def _update_file_count(self):
        """更新文件计数显示"""
        total_count = len(self.file_items)
        checked_count = sum(1 for info in self.file_items.values() if info["checked"].get())
        self.file_count_label.configure(text=f"已添加: {total_count} 个文件 | 已选中: {checked_count} 个")

    def _add_custom_keyword(self):
        """添加自定义关键词"""
        keyword = self.custom_kw_entry.get().strip()
        if not keyword:
            messagebox.showwarning("提示", "请输入关键词")
            return

        if keyword in self.custom_keywords:
            messagebox.showwarning("提示", "关键词已存在")
            return

        # 检查是否与默认关键词重复
        default_keywords_lower = [kw.lower() for kw in ["企密A", "企密AA", "企密AAA"]]
        if keyword.lower() in default_keywords_lower:
            messagebox.showwarning("提示", "该关键词与默认关键词重复")
            return

        self.custom_keywords.append(keyword)
        self._update_custom_keyword_list()
        self.custom_kw_entry.delete(0, 'end')
        self._log(f"已添加自定义关键词: {keyword}")

    def _delete_custom_keyword(self):
        """删除选中的自定义关键词"""
        if not self.custom_keywords:
            messagebox.showwarning("提示", "没有可删除的关键词")
            return

        # 创建选择对话框
        dialog = ctk.CTkToplevel(self)
        dialog.title("删除关键词")
        dialog.geometry("300x400")
        dialog.transient(self)
        dialog.grab_set()

        ctk.CTkLabel(
            dialog,
            text="选择要删除的关键词:",
            font=ctk.CTkFont(size=14, weight="bold")
        ).pack(pady=10)

        # 创建滚动框架
        scrollable_frame = ctk.CTkScrollableFrame(dialog, width=260, height=250)
        scrollable_frame.pack(pady=10, padx=20, fill="both", expand=True)

        # 复选框变量
        checkbox_vars = {}
        for kw in self.custom_keywords:
            var = ctk.BooleanVar()
            checkbox_vars[kw] = var
            ctk.CTkCheckBox(
                scrollable_frame,
                text=kw,
                variable=var
            ).pack(anchor="w", padx=10, pady=2)

        # 删除按钮
        def do_delete():
            to_delete = [kw for kw, var in checkbox_vars.items() if var.get()]
            if not to_delete:
                messagebox.showwarning("提示", "请至少选择一个关键词")
                return

            for kw in to_delete:
                self.custom_keywords.remove(kw)
                self._log(f"已删除自定义关键词: {kw}")

            self._update_custom_keyword_list()
            dialog.destroy()

        button_frame = ctk.CTkFrame(dialog)
        button_frame.pack(pady=10)

        ctk.CTkButton(
            button_frame,
            text="删除",
            command=do_delete,
            fg_color="red",
            width=100
        ).pack(side="left", padx=5)

        ctk.CTkButton(
            button_frame,
            text="取消",
            command=dialog.destroy,
            width=100
        ).pack(side="left", padx=5)

    def _update_custom_keyword_list(self):
        """更新自定义关键词列表显示"""
        self.custom_kw_listbox.delete("1.0", "end")
        for kw in self.custom_keywords:
            self.custom_kw_listbox.insert("end", kw + "\n")

    def _get_active_keywords(self) -> List[str]:
        """获取当前激活的所有关键词"""
        keywords = []

        # 默认关键词
        for kw, var in self.default_keyword_vars.items():
            if var.get():
                keywords.append(kw)

        # 自定义关键词
        keywords.extend(self.custom_keywords)

        return keywords

    def _start_check(self):
        """开始检查"""
        if self.is_checking:
            messagebox.showwarning("提示", "检查正在进行中...")
            return

        # 验证输入
        if not self.file_items:
            messagebox.showwarning("提示", "请先添加要检查的文件")
            return

        # 获取被勾选的文件
        selected_files = [path for path, info in self.file_items.items() if info["checked"].get()]

        if not selected_files:
            messagebox.showwarning("提示", "请至少勾选一个文件进行检查")
            return

        keywords = self._get_active_keywords()
        if not keywords:
            messagebox.showwarning("提示", "请至少选择一个关键词")
            return

        # 保存配置
        self.settings["default_keywords"] = {
            kw: var.get() for kw, var in self.default_keyword_vars.items()
        }
        self.settings["custom_keywords"] = self.custom_keywords.copy()
        self.settings["check_content"] = self.check_content_var.get()
        self.config_manager.save(self.settings)

        # 禁用按钮
        self._set_buttons_enabled(False)
        self.is_checking = True

        # 开始日志
        self.log_manager.start_logging()
        self._clear_log()
        self._log("=" * 50)
        self._log(f"开始检查，共 {len(selected_files)} 个文件")
        self._log(f"检查关键词: {', '.join(keywords)}")
        self._log(f"检测正文: {'是' if self.check_content_var.get() else '否'}")
        self._log("=" * 50)

        # 重置进度
        self.progress_bar.set(0)
        self.check_results = []

        # 启动检查线程
        threading.Thread(target=self._check_files_thread, args=(keywords, selected_files), daemon=True).start()

    def _check_files_thread(self, keywords: List[str], file_list: List[str]):
        """检查文件的线程函数"""
        total_files = len(file_list)
        total_matches = 0
        header_footer_matches = 0  # 页眉页脚命中数
        content_matches = 0  # 正文命中数
        failed_files = 0
        keyword_stats = {}  # 统计每个关键词的命中次数
        keyword_hf_stats = {}  # 页眉页脚关键词统计
        keyword_content_stats = {}  # 正文关键词统计

        checker = DocumentChecker(keywords, self.check_content_var.get())

        # 使用线程池
        with ThreadPoolExecutor(max_workers=4) as executor:
            future_to_file = {
                executor.submit(self._check_single_file, file_path, checker): file_path
                for file_path in file_list
            }

            completed = 0
            for future in as_completed(future_to_file):
                file_path = future_to_file[future]
                completed += 1

                try:
                    results = future.result()
                    file_name = os.path.basename(file_path)

                    if results:
                        # 分别统计页眉页脚和正文
                        hf_count = sum(1 for r in results if r.get('type') == 'header_footer')
                        content_count = sum(1 for r in results if r.get('type') == 'content')

                        total_matches += len(results)
                        header_footer_matches += hf_count
                        content_matches += content_count

                        # 显示命中信息
                        match_info = f"⚠️ [{completed}/{total_files}] {file_name} - 发现 {len(results)} 处命中"
                        if hf_count > 0 and content_count > 0:
                            match_info += f" (页眉页脚:{hf_count}, 正文:{content_count})"
                        elif hf_count > 0:
                            match_info += f" (页眉页脚:{hf_count})"
                        elif content_count > 0:
                            match_info += f" (正文:{content_count})"

                        self._log(match_info, "WARNING")

                        for result in results:
                            self._log(f"    → {result['keyword']} ({result['location']})")

                            # 统计关键词命中次数
                            kw = result['keyword']
                            result_type = result.get('type', 'header_footer')

                            keyword_stats[kw] = keyword_stats.get(kw, 0) + 1

                            if result_type == 'header_footer':
                                keyword_hf_stats[kw] = keyword_hf_stats.get(kw, 0) + 1
                            else:
                                keyword_content_stats[kw] = keyword_content_stats.get(kw, 0) + 1

                            self.check_results.append({
                                "file": file_path,
                                "keyword": result['keyword'],
                                "location": result['location'],
                                "context": result['context'],
                                "type": result_type
                            })
                    else:
                        self._log(f"✓ [{completed}/{total_files}] {file_name} - 未发现问题")

                except Exception as e:
                    failed_files += 1
                    self._log(f"✗ [{completed}/{total_files}] {os.path.basename(file_path)} - 检查失败: {str(e)}", "ERROR")

                # 更新进度
                progress = completed / total_files
                self.after(0, self.progress_bar.set, progress)
                self.after(0, self.status_label.configure, {"text": f"进度: {completed}/{total_files}"})

        # 完成 - 显示详细统计
        self._log("=" * 50)
        self._log(f"检查完成！")
        self._log(f"  总文件数: {total_files}")
        self._log(f"  成功检查: {total_files - failed_files}")
        self._log(f"  失败文件: {failed_files}")
        self._log(f"  总命中数: {total_matches}")
        self._log(f"    ├─ 页眉页脚: {header_footer_matches}")
        self._log(f"    └─ 正文内容: {content_matches}")

        if keyword_stats:
            self._log(f"\n关键词命中总统计:")
            for kw, count in sorted(keyword_stats.items(), key=lambda x: x[1], reverse=True):
                hf_count = keyword_hf_stats.get(kw, 0)
                c_count = keyword_content_stats.get(kw, 0)
                self._log(f"  {kw}: {count} 次 (页眉页脚:{hf_count}, 正文:{c_count})")

        self._log("=" * 50)

        self.after(0, self._check_complete)

    def _check_single_file(self, file_path: str, checker: DocumentChecker) -> List[Dict]:
        """检查单个文件"""
        ext = Path(file_path).suffix.lower()

        if ext == '.docx':
            return checker.check_word(file_path)
        elif ext == '.xlsx':
            return checker.check_excel(file_path)
        elif ext == '.pptx':
            return checker.check_ppt(file_path)
        else:
            return []

    def _check_complete(self):
        """检查完成回调"""
        self.is_checking = False
        self._set_buttons_enabled(True)
        self.status_label.configure(text="检查完成")
        self.progress_bar.set(1.0)

        if self.check_results:
            messagebox.showinfo(
                "检查完成",
                f"检查完成！\n共发现 {len(self.check_results)} 处关键词命中\n\n请查看日志获取详细信息"
            )
        else:
            messagebox.showinfo("检查完成", "检查完成！未发现敏感关键词")

    def _set_buttons_enabled(self, enabled: bool):
        """设置按钮启用状态"""
        state = "normal" if enabled else "disabled"
        self.start_btn.configure(state=state)
        self.select_files_btn.configure(state=state)
        self.select_folder_btn.configure(state=state)

    def _log(self, message: str, level="INFO"):
        """添加日志"""
        log_entry = self.log_manager.log(message, level)
        self.log_text.insert("end", log_entry + "\n")
        self.log_text.see("end")

    def _clear_log(self):
        """清空日志"""
        self.log_text.delete("1.0", "end")

    def _save_log(self):
        """保存日志"""
        file_path = filedialog.asksaveasfilename(
            title="保存日志",
            defaultextension=".txt",
            filetypes=[("文本文件", "*.txt"), ("所有文件", "*.*")]
        )
        if file_path:
            if self.log_manager.save_to_file(file_path):
                messagebox.showinfo("成功", "日志已保存")
            else:
                messagebox.showerror("错误", "日志保存失败")

    def _export_report(self):
        """导出Excel报告"""
        if not self.check_results:
            messagebox.showwarning("提示", "没有可导出的检查结果")
            return

        file_path = filedialog.asksaveasfilename(
            title="导出报告",
            defaultextension=".xlsx",
            filetypes=[("Excel文件", "*.xlsx"), ("所有文件", "*.*")]
        )

        if file_path:
            try:
                from openpyxl import Workbook

                wb = Workbook()
                ws = wb.active
                ws.title = "检查报告"

                # 表头
                headers = ["文件路径", "关键词", "位置", "类型", "上下文"]
                ws.append(headers)

                # 数据
                for result in self.check_results:
                    type_text = "页眉页脚" if result.get('type') == 'header_footer' else "正文"
                    ws.append([
                        result['file'],
                        result['keyword'],
                        result['location'],
                        type_text,
                        result['context']
                    ])

                # 调整列宽
                ws.column_dimensions['A'].width = 50
                ws.column_dimensions['B'].width = 15
                ws.column_dimensions['C'].width = 30
                ws.column_dimensions['D'].width = 12
                ws.column_dimensions['E'].width = 50

                wb.save(file_path)
                messagebox.showinfo("成功", f"报告已导出到:\n{file_path}")

            except Exception as e:
                messagebox.showerror("错误", f"导出报告失败:\n{str(e)}")

    def on_closing(self):
        """窗口关闭事件"""
        if self.is_checking:
            if messagebox.askokcancel("确认", "检查正在进行中，确定要退出吗?"):
                self.destroy()
        else:
            self.destroy()


# ==================== 主函数 ====================
def main():
    """主函数"""
    app = DocumentCheckerGUI()
    app.protocol("WM_DELETE_WINDOW", app.on_closing)
    app.mainloop()


if __name__ == "__main__":
    main()
